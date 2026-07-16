# -*- coding: utf-8 -*-
"""Build the RAG×KG technical briefing deck (Traditional Chinese, 16:9)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

SHOTS = r"C:\Users\G635LXG\Downloads\RAG\AI_Document_V3\.shots"
DECK = os.path.join(SHOTS, "deck")
OUT = r"C:\Users\G635LXG\Downloads\RAG\RAG_KG_技術匯報.pptx"
FONT = "Microsoft JhengHei"

# palette
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x16, 0x77, 0xFF)
LBLUE = RGBColor(0xE7, 0xF0, 0xFF)
BG = RGBColor(0xF5, 0xF7, 0xFA)
RED = RGBColor(0xE0, 0x31, 0x31)
GREEN = RGBColor(0x2F, 0x9E, 0x44)
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=DARK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # ensure east-asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=2):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if hasattr(para, 'get') is False and isinstance(para, dict):
            pass
        for seg in para:
            txt, sz, bd, col = (seg + (None,) * 4)[:4]
            r = p.add_run(); r.text = txt
            _set_font(r, sz, bd or False, col or DARK)
    return tb


def title_bar(s, kicker, title):
    box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    box(s, 0, Inches(1.15), SW, Pt(3), fill=BLUE)
    if kicker:
        text(s, Inches(0.6), Inches(0.16), Inches(11), Inches(0.3), [[(kicker, 12, True, RGBColor(0x9E,0xC5,0xFF))]])
    text(s, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.6), [[(title, 26, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, items, l=Inches(0.7), t=Inches(1.5), w=Inches(12), h=Inches(5.6), sz=16, gap=8):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, txt, col, bold = (it + (0, "", DARK, False))[:4] if isinstance(it, tuple) else (0, it, DARK, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = lvl
        mark = "●  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark; _set_font(r, sz-2, False, BLUE if lvl == 0 else GRAY)
        r2 = p.add_run(); r2.text = txt; _set_font(r2, sz if lvl == 0 else sz-2, bold, col)
    return tb


def img_size(path):
    with Image.open(path) as im:
        return im.size


def place_img(s, path, bl, bt, bw, bh, border=True):
    iw, ih = img_size(path)
    r = min(bw / iw, bh / ih)
    w, h = int(iw * r), int(ih * r)
    l = bl + (bw - w) // 2; t = bt + (bh - h) // 2
    if border:
        box(s, l-Emu(9000), t-Emu(9000), w+Emu(18000), h+Emu(18000), fill=WHITE, line=LGRAY, line_w=1)
    s.shapes.add_picture(path, l, t, w, h)
    return l, t, w, h


def highlight(s, il, it, iw, ih, xf, yf, wf, hf, label="", lside="right"):
    """red rounded box over the image region (fractional coords) + optional callout."""
    l = il + int(iw * xf); t = it + int(ih * yf); w = int(iw * wf); h = int(ih * hf)
    hb = box(s, l, t, w, h, fill=None, line=RED, line_w=2.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        hb.adjustments[0] = 0.08
    except Exception:
        pass
    if label:
        tw, thh = Inches(1.85), Inches(0.32)
        if lside == "right":
            cl = l + w + Emu(45000); ct = t + h//2 - thh//2
        elif lside == "left":
            cl = max(Inches(0.05), l - tw - Emu(45000)); ct = t + h//2 - thh//2
        elif lside == "top":
            cl = l + w//2 - tw//2; ct = max(Inches(0.05), t - thh - Emu(30000))
        else:
            cl = l + w//2 - tw//2; ct = t + h + Emu(30000)
        tag = box(s, cl, ct, tw, thh, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = tag.text_frame; tf.word_wrap = True; tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
        tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = label; _set_font(rr, 10.5, True, WHITE)


def chip(s, l, t, w, h, label, fill, fg=WHITE, sz=12):
    c = box(s, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left=Pt(3); tf.margin_right=Pt(3); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, sz, True, fg)
    return c


def arrow(s, l, t, w, h, color=BLUE):
    a = box(s, l, t, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)
    return a


# ============================ SLIDES ============================

# 1 cover
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(4.55), SW, Pt(4), fill=BLUE)
box(s, Inches(0.0), 0, Inches(0.22), SH, fill=BLUE)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4), [[("AI 文件智慧問答系統 — 技術匯報", 16, True, RGBColor(0x9E,0xC5,0xFF))]])
text(s, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.6),
     [[("從 RAG 到 RAG × KG", 46, True, WHITE)], [("測試規範智慧問答的準確性強化之路", 30, True, RGBColor(0xCF,0xE2,0xFF))]])
text(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.2),
     [[("以 MIL-STD-810H 軍規環境測試文件為例", 18, False, RGBColor(0xCF,0xDB,0xEA))],
      [("檢索增強生成（RAG）× 知識圖譜（KG）× 混合查詢路由", 15, False, GRAY)]])

# 2 agenda
s = slide(); title_bar(s, "AGENDA", "本日大綱")
ag = [
    "一、RAG 基礎：什麼是檢索增強生成、為什麼需要它",
    "二、系統介紹：主介面、核心功能、系統設定",
    "三、基本 RAG 會踩的坑（六大陷阱）",
    "四、本專案對 RAG 查詢準確性的五大強化",
    "五、RAG 的致命傷：跨文件關聯",
    "六、知識圖譜 KG：是什麼、為何對測試規範特別有效",
    "七、RAG × KG 結合：架構與運作",
    "八、AS-IS vs TO-BE 效果比較",
    "九、量化成果與總結",
]
bullets(s, [(0, a, DARK, False) for a in ag], t=Inches(1.55), sz=19, gap=12)


def section(num, title, sub=""):
    s = slide(); box(s, 0, 0, SW, SH, fill=NAVY); box(s, Inches(0.9), Inches(3.05), Inches(1.2), Pt(5), fill=BLUE)
    text(s, Inches(0.9), Inches(2.1), Inches(2), Inches(1), [[(num, 60, True, RGBColor(0x2E,0x5C,0x9E))]])
    text(s, Inches(0.9), Inches(3.3), Inches(11), Inches(1), [[(title, 40, True, WHITE)]])
    if sub:
        text(s, Inches(0.95), Inches(4.5), Inches(11), Inches(0.6), [[(sub, 18, False, RGBColor(0xAF,0xC6,0xE0))]])
    return s

# ---- Part 1 RAG basics ----
section("01", "RAG 基礎", "檢索增強生成 — 讓 LLM「查了再答」")

# 4 what is RAG (flow)
s = slide(); title_bar(s, "PART 1 · RAG 基礎", "什麼是 RAG（檢索增強生成）")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
     [[("RAG = Retrieval-Augmented Generation。先「檢索」相關文件片段，再讓 LLM 依據這些片段「生成」帶出處的答案，而非憑模型記憶硬答。", 15, False, DARK)]])
steps = ["文件\nPDF/規範", "切塊\nChunk", "向量化\nEmbedding", "向量庫\nFAISS", "提問檢索\nTop-K", "LLM 生成\n帶引用答案"]
n = len(steps); x = Inches(0.55); y = Inches(2.7); bw = Inches(1.72); bh = Inches(1.25); gap = Inches(0.32)
cols = [NAVY, BLUE, BLUE, RGBColor(0x37,0x99,0x5B), ORANGE, RGBColor(0x7A,0x3E,0xA8)]
for i, st in enumerate(steps):
    bx = x + i * (bw + gap)
    c = box(s, bx, y, bw, bh, fill=cols[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for j, ln in enumerate(st.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, 14 if j == 0 else 11, j == 0, WHITE)
    if i < n - 1:
        arrow(s, bx + bw + Emu(20000), y + bh//2 - Inches(0.13), gap - Emu(40000), Inches(0.26), color=GRAY)
text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
     [[("關鍵：答案來自「檢索到的真實文件」，每一句都能附上來源頁碼 → 可溯源、可查證。", 15, True, BLUE)]])
box(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.4), fill=BG, line=LGRAY)
text(s, Inches(0.95), Inches(5.35), Inches(11.4), Inches(1.1),
     [[("本系統的 RAG 管線：", 14, True, NAVY)],
      [("pdfminer 抽文字 → 切塊 → qwen3-embedding:8b 向量化 → FAISS 向量庫 → 向量檢索 + Cross-encoder 重排 → gemma4:12b 生成帶引用答案", 13, False, DARK)]])

# 5 why RAG
s = slide(); title_bar(s, "PART 1 · RAG 基礎", "為什麼需要 RAG")
cards = [("① 消除幻覺", "LLM 不再憑記憶亂編，答案必須有檢索到的文件支撐。", BLUE),
         ("② 私有/時效知識", "軍規測試文件無法內建於模型，RAG 讓模型即時讀取私有語料。", GREEN),
         ("③ 可溯源", "每個答案附來源頁碼，使用者可點開 PDF 查證 — 合規場景的剛需。", ORANGE)]
cw = Inches(3.95); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(1.7)
for i, (h, b, col) in enumerate(cards):
    cx = x0 + i * (cw + gap)
    box(s, cx, y0, cw, Inches(3.6), fill=WHITE, line=LGRAY, line_w=1.2)
    box(s, cx, y0, cw, Inches(0.85), fill=col)
    text(s, cx+Inches(0.25), y0+Inches(0.18), cw-Inches(0.5), Inches(0.5), [[(h, 19, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+Inches(0.28), y0+Inches(1.15), cw-Inches(0.56), Inches(2.2), [[(b, 15, False, DARK)]])
text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
     [[("→ 對「測試規範查詢」這種需要精確、可查證、來自私有大量 PDF 的場景，RAG 是基礎門檻。", 16, True, NAVY)]])

# ---- Part 2 system intro ----
section("02", "系統介紹", "主介面 · 核心功能 · 系統設定")


def img_slide(kicker, title, img, caption, hl=None, box_t=Inches(1.45), box_h=Inches(5.0)):
    s = slide(); title_bar(s, kicker, title)
    il, it, iw, ih = place_img(s, img, Inches(0.7), box_t, Inches(11.9), box_h)
    if hl:
        for spec in hl:
            highlight(s, il, it, iw, ih, *spec)
    if caption:
        text(s, Inches(0.7), Inches(6.62), Inches(11.9), Inches(0.6), [[(caption, 13, False, GRAY)]], align=PP_ALIGN.CENTER)
    return s

def marker(s, il, it, iw, ih, xf, yf, n, d=Inches(0.36)):
    cx = il + int(iw * xf) - d//2; cy = it + int(ih * yf) - d//2
    mk = box(s, cx, cy, d, d, fill=RED, line=WHITE, line_w=1.5, shape=MSO_SHAPE.OVAL)
    tf = mk.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); _set_font(r, 14, True, WHITE)

# functions overview (grid first)
s = slide(); title_bar(s, "PART 2 · 系統介紹", "功能總覽：一套系統涵蓋全流程")
funcs = [("① 文件列表", "瀏覽/檢視/下載已匯入的規範文件", BLUE), ("② 建立文件", "上傳 PDF、OCR / Vision 解析、向量化", GREEN),
         ("③ RAG 智慧問答", "三模式：純RAG / 混合 / Agent", ORANGE), ("④ 知識圖譜", "規範引用關係、版本鏈視覺化", RGBColor(0x7A,0x3E,0xA8)),
         ("⑤ 我的筆記本", "儲存查詢結果與重點", RGBColor(0x0E,0x7A,0x8B)), ("⑥ 管理介面", "模型 / OCR / 向量化 / 提示詞設定", NAVY),
         ("⑦ 向量查詢測試", "直接測試向量檢索品質", RGBColor(0xC2,0x41,0x0C)), ("⑧ 向量庫健康", "FAISS 與 DB 一致性檢查", RGBColor(0x2B,0x8A,0x3E))]
x0 = Inches(0.7); y0 = Inches(1.7); cw = Inches(3.0); ch = Inches(1.55); gx = Inches(0.18); gy = Inches(0.2)
for i, (h, d, col) in enumerate(funcs):
    r, c = divmod(i, 4)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    box(s, cx, cy, cw, ch, fill=WHITE, line=LGRAY)
    box(s, cx, cy, Pt(6), ch, fill=col)
    text(s, cx+Inches(0.22), cy+Inches(0.16), cw-Inches(0.4), Inches(0.4), [[(h, 15, True, col)]])
    text(s, cx+Inches(0.22), cy+Inches(0.72), cw-Inches(0.4), Inches(0.8), [[(d, 12, False, DARK)]])
text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.6),
     [[("流程：①② 匯入解析 → ③④ 問答與圖譜 → ⑤ 收藏 → ⑥⑦⑧ 設定與診斷。以下逐一介紹。", 15, True, NAVY)]])

# main interface — feature-by-feature (numbered markers + legend)
s = slide(); title_bar(s, "PART 2 · 系統介紹", "③ 主介面詳解：RAG 智慧問答")
il, it, iw, ih = place_img(s, os.path.join(DECK, "02_qa_main.png"), Inches(0.45), Inches(1.5), Inches(7.6), Inches(5.0))
feat = [(0.30, 0.083, "三模式切換", "純RAG / 混合（自動路由）/ Agent"),
        (0.225, 0.305, "問題輸入", "自然語言提問，Ctrl+Enter 送出"),
        (0.225, 0.44, "範圍篩選", "可限定分類 / 專案 / 資料夾 / 特定文件"),
        (0.225, 0.635, "來源筆數", "控制檢索回傳的段落數量 (Top-K)"),
        (0.225, 0.705, "送出查詢", "送出後右側即時串流答案"),
        (0.60, 0.115, "對話結果區", "答案 + 可點來源 + 追問")]
for i, (xf, yf, _, _) in enumerate(feat, 1):
    marker(s, il, it, iw, ih, xf, yf, i)
# legend
lx = Inches(8.35); ly = Inches(1.6)
box(s, lx, ly, Inches(4.55), Inches(4.9), fill=BG, line=LGRAY)
tb = s.shapes.add_textbox(lx+Inches(0.2), ly+Inches(0.18), Inches(4.15), Inches(4.6)); tf = tb.text_frame; tf.word_wrap = True
for i, (_, _, h, d) in enumerate(feat, 1):
    p = tf.paragraphs[0] if i == 1 else tf.add_paragraph(); p.space_after = Pt(9)
    r = p.add_run(); r.text = f"{i}  "; _set_font(r, 14, True, RED)
    r2 = p.add_run(); r2.text = h; _set_font(r2, 14, True, NAVY)
    p2 = tf.add_paragraph(); p2.space_after = Pt(9)
    r3 = p2.add_run(); r3.text = "      " + d; _set_font(r3, 12, False, DARK)

# documents + new
s = slide(); title_bar(s, "PART 2 · 系統介紹", "①② 文件管理：列表與匯入")
place_img(s, os.path.join(DECK, "05_documents.png"), Inches(0.5), Inches(1.5), Inches(6.1), Inches(4.7))
place_img(s, os.path.join(DECK, "06_doc_new.png"), Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.7))
text(s, Inches(0.5), Inches(6.35), Inches(6.1), Inches(0.9),
     [[("① 文件列表", 14, True, BLUE)], [("已匯入 40 份軍規文件，可檢視 / 下載 / 看分類結果。", 12.5, False, DARK)]], align=PP_ALIGN.CENTER)
text(s, Inches(6.8), Inches(6.35), Inches(6.0), Inches(0.9),
     [[("② 建立文件", 14, True, GREEN)], [("上傳 PDF → 可選強制 OCR / Vision 解析 → 自動向量化。", 12.5, False, DARK)]], align=PP_ALIGN.CENTER)

# KG page intro
s = slide(); title_bar(s, "PART 2 · 系統介紹", "④ 知識圖譜頁")
place_img(s, os.path.join(DECK, "03_kg_full.png"), Inches(0.7), Inches(1.55), Inches(7.4), Inches(4.9))
box(s, Inches(8.4), Inches(1.7), Inches(4.5), Inches(4.5), fill=BG, line=LGRAY)
bullets(s, [(0, "視覺化規範之間的引用網與版本鏈", DARK, False),
            (0, "左側：實體 / 關係統計與型別分布", DARK, False),
            (0, "可搜尋任一規範為中心展開子圖", DARK, False),
            (0, "點節點放大可見名稱與關係", DARK, False),
            (0, "（KG 的原理與效果見第六、七部分）", GRAY, False)],
        l=Inches(8.6), t=Inches(1.95), w=Inches(4.1), h=Inches(4.2), sz=14, gap=12)

# settings (annotated)
img_slide("PART 2 · 系統介紹", "⑥ 系統設定：模型、OCR、向量化",
          os.path.join(DECK, "09b_settings.png"),
          "管理介面 → 系統設置：可線上切換 LLM / Embedding / Vision 模型、OCR 引擎、切塊與檢索參數。",
          hl=[(0.03, 0.015, 0.94, 0.05, "系統概覽", "right"),
              (0.03, 0.085, 0.94, 0.125, "模型設定", "right"),
              (0.03, 0.30, 0.94, 0.075, "OCR 引擎", "right"),
              (0.03, 0.62, 0.94, 0.10, "向量化參數", "right")])

# tools (4-up incl notebook)
s = slide(); title_bar(s, "PART 2 · 系統介紹", "⑤⑦⑧ 工具：筆記本 / 向量查詢 / 向量健康 / 提示詞")
imgs = [("10_notebook.png", "⑤ 我的筆記本", "收藏查詢結果與重點"),
        ("07_vector_search.png", "⑦ 向量查詢測試", "直接驗證檢索品質"),
        ("08_vector_health.png", "⑧ 向量庫健康", "FAISS ↔ DB 一致性"),
        ("09c_prompts.png", "⑥ RAG 提示詞", "可調整回答風格")]
x0 = Inches(0.35); cw = Inches(3.18); gx = Inches(0.08)
for i, (fn, h, d) in enumerate(imgs):
    cx = x0 + i*(cw+gx)
    place_img(s, os.path.join(DECK, fn), cx, Inches(1.55), cw, Inches(3.9))
    text(s, cx, Inches(5.65), cw, Inches(0.4), [[(h, 13.5, True, NAVY)]], align=PP_ALIGN.CENTER)
    text(s, cx, Inches(6.05), cw, Inches(0.5), [[(d, 12, False, GRAY)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.35), Inches(6.6), Inches(12.6), Inches(0.5), [[("可觀測、可調參、可診斷 — 不是黑盒子。", 14, True, NAVY)]], align=PP_ALIGN.CENTER)

# ---- Part 3 pitfalls ----
section("03", "基本 RAG 會踩的坑", "為什麼「能跑」不等於「準確」")
s = slide(); title_bar(s, "PART 3 · RAG 的坑", "基本 RAG 的六大陷阱")
pits = [("① 切塊破壞語意", "關鍵句、表格被切斷，檢索到殘缺片段。"),
        ("② 檢索漏抓", "用詞不同但語意相近 → Top-K 沒撈到；「引用清單」這種列表最常漏。"),
        ("③ 幻覺", "檢索不足時模型硬掰答案，且看起來很有自信。"),
        ("④ 思考模型暴走/超時", "reasoning token 爆量，單次查詢動輒逾時。"),
        ("⑤ Lost-in-the-middle", "context 塞太多，中段內容被模型忽略。"),
        ("⑥ 跨文件無能", "只在單一片段內找答案，無法串接『A 引用 B、B 取代 C』。")]
x0 = Inches(0.7); y0 = Inches(1.6); cw = Inches(5.85); ch = Inches(1.5); gx = Inches(0.3); gy = Inches(0.25)
for i, (h, d) in enumerate(pits):
    r, c = divmod(i, 2)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    box(s, cx, cy, cw, ch, fill=RGBColor(0xFB,0xE9,0xE7), line=RGBColor(0xF3,0xB7,0xA8))
    text(s, cx+Inches(0.25), cy+Inches(0.14), cw-Inches(0.5), Inches(0.45), [[(h, 16, True, RGBColor(0xC0,0x39,0x2B))]])
    text(s, cx+Inches(0.25), cy+Inches(0.62), cw-Inches(0.5), Inches(0.8), [[(d, 13.5, False, DARK)]])
text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
     [[("⑥ 跨文件無能 ＝ 本報告後半的主角，也是純 RAG 的天花板。", 14, True, ORANGE)]])

# ---- Part 4 enhancements ----
section("04", "RAG 準確性的五大強化", "把「能跑」調成「準、快、不亂編」")

# overview
s = slide(); title_bar(s, "PART 4 · 準確性強化", "強化總覽：檢索 → 重排 → 擴展 → 合成 → 路由")
pipe = [("① 重排序", "Cross-encoder\n精排相關性", BLUE), ("② 小找大", "鄰塊擴展\n補回上下文", GREEN),
        ("③ 關閉思考", "think=false\n解暴走超時", ORANGE), ("④ Grounded 合成", "強制引用\n誠實兜底", RGBColor(0x7A,0x3E,0xA8)),
        ("⑤ 混合路由", "依題型\n選最佳路徑", RGBColor(0x0E,0x7A,0x8B))]
x = Inches(0.5); y = Inches(2.4); bw = Inches(2.3); bh = Inches(1.7); gap = Inches(0.18)
for i, (h, d, col) in enumerate(pipe):
    bx = x + i*(bw+gap)
    box(s, bx, y, bw, bh, fill=WHITE, line=col, line_w=1.6)
    box(s, bx, y, bw, Inches(0.6), fill=col)
    text(s, bx, y+Inches(0.1), bw, Inches(0.4), [[(h, 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx+Inches(0.1), y+Inches(0.75), bw-Inches(0.2), Inches(0.9),
         [[(ln, 12.5, False, DARK)] for ln in d.split("\n")], align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, bx+bw+Emu(5000), y+bh//2-Inches(0.1), gap-Emu(10000), Inches(0.2), color=GRAY)
text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.5),
     [[("五道工序層層把關，讓答案「找得到（②）、排得對（①）、不超時（③）、不亂編（④）、走對路（⑤）」。", 15, True, NAVY)]])
box(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.3), fill=BG, line=LGRAY)
text(s, Inches(0.95), Inches(5.45), Inches(11.4), Inches(1.0),
     [[("成果預覽：RAG 由「思考模型動輒逾時」降到 10–27 秒；內容題能回完整數據表、表格、程序；證據不足時老實說查不到（零幻覺）。", 14, False, DARK)]])


def enh_slide(idx, title, what, how, color, example, img=None, hl=None):
    s = slide(); title_bar(s, f"PART 4 · 強化{idx}", title)
    lw = Inches(5.5) if img else Inches(11.9)
    box(s, Inches(0.7), Inches(1.5), lw, Inches(1.25), fill=LBLUE, line=color, line_w=1.2)
    text(s, Inches(0.95), Inches(1.6), lw-Inches(0.5), Inches(0.35), [[("作用", 13.5, True, color)]])
    text(s, Inches(0.95), Inches(1.98), lw-Inches(0.5), Inches(0.75), [[(what, 14, False, DARK)]])
    box(s, Inches(0.7), Inches(2.9), lw, Inches(2.0), fill=WHITE, line=LGRAY)
    text(s, Inches(0.95), Inches(3.0), lw-Inches(0.5), Inches(0.35), [[("怎麼做", 13.5, True, NAVY)]])
    bullets(s, [(0, h, DARK, False) for h in how], l=Inches(0.95), t=Inches(3.4), w=lw-Inches(0.5), h=Inches(1.5), sz=13, gap=6)
    # example box (light yellow)
    box(s, Inches(0.7), Inches(5.05), lw, Inches(1.55), fill=RGBColor(0xFF,0xF7,0xE0), line=RGBColor(0xF0,0xC3,0x6B), line_w=1.2)
    text(s, Inches(0.95), Inches(5.15), lw-Inches(0.5), Inches(0.35), [[("💡 簡單案例", 13.5, True, RGBColor(0xB5,0x7A,0x0A))]])
    text(s, Inches(0.95), Inches(5.55), lw-Inches(0.5), Inches(1.0), [[(example, 13.5, False, DARK)]])
    if img:
        il, it, iw, ih = place_img(s, img, Inches(6.4), Inches(1.55), Inches(6.4), Inches(4.9))
        if hl:
            for spec in hl:
                highlight(s, il, it, iw, ih, *spec)
    return s

enh_slide("①", "Cross-encoder 重排序（Rerank）",
          "向量初篩 Top-N 後，用交叉編碼器對「問題×段落」逐一精算相關性，把真正相關的拉到最前面，解決「檢索排序不準 / 漏抓」。",
          ["向量檢索快但粗：只比對「整體語意相似」", "Cross-encoder（bge-reranker）讓問題與每個段落『一起進模型』精算分數",
           "重排後 Top-K 命中率大幅提升，雜訊段落被壓下去"], BLUE,
          "問「鹽霧鹽水濃度多少？」→ 向量初篩可能把『鹽霧腐蝕原理』排在前面；重排後把真正含『5±1% 濃度表』的段落拉到第 1。")

enh_slide("②", "小找大 chunk 擴展（Small-to-Big）",
          "命中小塊後，自動把前後鄰近段落補回來，讓 LLM 拿到完整上下文，解決「切塊破壞語意」。",
          ["檢索用「小塊」→ 精準命中", "回答用「大塊」→ 命中塊 + 鄰塊一起餵給 LLM",
           "表格、跨段落的步驟不再被切斷"], GREEN,
          "命中『鹽濃度 5±1%』這一小塊 → 自動補回前後段（量測溫度 35±1°C、比重 1.025–1.037），答案才完整、不殘缺。")

enh_slide("③", "關閉思考（think = false）",
          "讓 gemma / qwen 等思考型模型停止暴走 reasoning，直接給答案，解決「超時 / token 爆量」。",
          ["思考模型預設輸出大量 reasoning token → 慢且常逾時", "Ollama 請求顯式送 think=false（並修正 RAG 內部 hardcoded think=True）",
           "RAG 由「動輒逾時」降到 10–27 秒"], ORANGE,
          "同一題：gemma4 思考版單次生成 18.8 秒（常逾時）→ 關閉思考後 1.9 秒。速度差近 10 倍。")

enh_slide("④", "Grounded 合成 + 強制引用 + 誠實兜底",
          "答案每段掛 [來源N]；證據不足時老實說「查不到」，絕不亂編，解決「幻覺」。",
          ["合成階段強制標註來源，前端可點開對應 PDF 頁", "信心不足 → 回『此資訊不在知識庫』而非硬掰"],
          RGBColor(0x7A,0x3E,0xA8),
          "問『低壓測試模擬到多高？』→ 文件其實寫『由測試計畫定義』。系統不編一個數字，而是誠實回答『無固定值，由測試計畫指定』。",
          img=os.path.join(SHOTS, "src_answer.png"),
          hl=[(0.40, 0.78, 0.55, 0.16, "每筆來源可點預覽", "top")])

enh_slide("⑤", "混合查詢路由（自動判定題型）",
          "自動分類問題：內容題走最快最完整的純 RAG，關係題走 KG+Agent，避免繞冤枉路也避免淡化答案。",
          ["內容題（怎麼進行 / 數值）→ 純 RAG（快、完整）", "關係題（引用 / 版本 / 有哪些）→ KG + Agent"],
          RGBColor(0x0E,0x7A,0x8B),
          "『鹽霧測試怎麼進行？』→ 判為內容題 → 純 RAG。\n『510.7 引用了哪些規範？』→ 判為關係題 → KG + Agent。",
          img=os.path.join(DECK, "02_qa_main.png"),
          hl=[(0.305, 0.085, 0.085, 0.028, "三模式", "top")])

# ---- Part 5 fatal weakness ----
section("05", "RAG 的致命傷", "調到最佳，跨文件關聯仍然無解")
s = slide(); title_bar(s, "PART 5 · 致命傷", "即使 RAG 調到最佳…跨文件關聯仍無解")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
     [[("RAG 看的是「文字相似」，看不到「規範之間的明確關係」。以下兩題，純 RAG 都答不好：", 15, True, DARK)]])
ex = [("Q：Method 510.7 引用了哪些外部規範？",
       "純 RAG：只撈到模糊片段，漏掉完整的 reference 清單（列表型資訊最容易漏抓）。"),
      ("Q：MIL-HDBK-310 和 MIL-STD-210 是什麼關係？",
       "純 RAG：直接回「查無相關資料」— 因為關係不寫在同一段文字裡，向量檢索看不到。"),
      ("Q：MIL-STD-810 有哪些版本演進？",
       "純 RAG：只認得當前版，給不出 B→C→…→H 的完整版本鏈。")]
y = Inches(2.1)
for q, a in ex:
    box(s, Inches(0.7), y, Inches(11.9), Inches(1.35), fill=WHITE, line=LGRAY)
    box(s, Inches(0.7), y, Pt(6), Inches(1.35), fill=ORANGE)
    text(s, Inches(0.95), y+Inches(0.13), Inches(11.4), Inches(0.45), [[(q, 15, True, NAVY)]])
    text(s, Inches(0.95), y+Inches(0.62), Inches(11.4), Inches(0.7), [[("✗ ", 14, True, RED), (a, 14, False, DARK)]])
    y += Inches(1.5)
text(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.5),
     [[("規範世界本質是一張「引用網 + 版本鏈」→ 這正是「圖（Graph）」的問題，需要知識圖譜 KG。", 15, True, ORANGE)]])

# ---- Part 6 KG ----
section("06", "知識圖譜 KG", "把「規範之間的關係」結構化")

# what is KG (concept diagram)
s = slide(); title_bar(s, "PART 6 · 知識圖譜", "什麼是知識圖譜（KG）")
text(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
     [[("KG = 用「實體（節點）」+「關係（邊）」把知識結構化。電腦因此能『沿著關係走』，而非只比對文字。", 15, False, DARK)]])
# nodes
def node(cx, cy, label, col, w=Inches(2.3), h=Inches(0.8)):
    nb = box(s, cx, cy, w, h, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = nb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, 14, True, WHITE)
    return cx, cy, w, h
def edge(x1, y1, x2, y2, label, col=GRAY):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = col; cn.line.width = Pt(1.75)
    mx, my = (x1+x2)//2 - Inches(0.9), (y1+y2)//2 - Inches(0.18)
    tag = box(s, mx, my, Inches(1.8), Inches(0.34), fill=WHITE, line=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = label; _set_font(rr, 11, True, col)
n1 = node(Inches(5.2), Inches(2.35), "MIL-STD-810H", NAVY)
n2 = node(Inches(1.4), Inches(4.4), "ASTM B117", RGBColor(0x37,0x99,0x5B))
n3 = node(Inches(5.2), Inches(4.9), "MIL-STD-810G", ORANGE)
n4 = node(Inches(9.0), Inches(4.4), "MIL-HDBK-310", RGBColor(0x7A,0x3E,0xA8))
edge(Inches(5.6), Inches(3.15), Inches(2.5), Inches(4.4), "references 引用", RGBColor(0x37,0x99,0x5B))
edge(Inches(6.35), Inches(3.15), Inches(6.35), Inches(4.9), "supersedes 取代", ORANGE)
edge(Inches(6.9), Inches(3.15), Inches(9.6), Inches(4.4), "references 引用", RGBColor(0x7A,0x3E,0xA8))
text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.8),
     [[("封閉關係詞彙：references（引用）/ supersedes（取代）/ defines（定義）/ requires（要求）/ derives_from（衍生）。", 14, True, BLUE)]])

# real KG (overview + zoomed-with-labels)
s = slide(); title_bar(s, "PART 6 · 知識圖譜", "本系統的 KG 長什麼樣")
il, it, iw, ih = place_img(s, os.path.join(DECK, "03_kg_full.png"), Inches(0.5), Inches(1.5), Inches(6.2), Inches(4.4))
highlight(s, il, it, iw, ih, 0.11, 0.095, 0.165, 0.36, "規模 + 型別分布", "right")
place_img(s, os.path.join(DECK, "kg_clean_7808.png"), Inches(6.95), Inches(1.5), Inches(5.95), Inches(4.4))
text(s, Inches(0.5), Inches(6.0), Inches(6.2), Inches(1.0),
     [[("全圖：3,074 實體 / 4,979 關係；節點型別含 規範 / 方法 / 章節 / 文件。", 13, False, GRAY)]], align=PP_ALIGN.CENTER)
text(s, Inches(6.95), Inches(6.0), Inches(5.95), Inches(1.0),
     [[("放大後可見節點文字：以 MIL-PRF-7808 為中心 — 連到版本 7808L、章節 2.2.2、被 810H 引用。", 13, False, GRAY)]], align=PP_ALIGN.CENTER)

# why KG effective
s = slide(); title_bar(s, "PART 6 · 知識圖譜", "為什麼 KG 對「測試規範跨文件查詢」特別有效")
rows = [("規範世界 = 引用網 + 版本鏈", "ISO/MIL-STD/ASTM 互相引用、版本不斷演進 — 天生就是圖結構。"),
        ("確定性、不漏項", "「方法 510.7 引用哪些規範」直接讀圖上的邊，完整且精確（RAG 會漏列表）。"),
        ("跨文件多跳", "可走『810H → 引用 → MIL-HDBK-310 → 又引用 → MIL-STD-210』的引用鏈。"),
        ("章節/方法級細化", "精確到『§4.1.1.4 引用了 ASTM D185-07』，而非只到文件層級。"),
        ("反向與版本查詢", "「誰引用了 X」「X 被什麼取代」這類 RAG 答不出的關係，KG 一查即得。")]
y = Inches(1.65)
for h, d in rows:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.92), fill=WHITE, line=LGRAY)
    box(s, Inches(0.7), y, Inches(3.8), Inches(0.92), fill=LBLUE)
    text(s, Inches(0.9), y+Inches(0.06), Inches(3.5), Inches(0.8), [[(h, 14.5, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.75), y+Inches(0.06), Inches(7.6), Inches(0.8), [[(d, 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.04)

# ---- Part 7 RAG x KG ----
section("07", "RAG × KG 結合", "KG 給骨架，RAG 給血肉")

# architecture
s = slide(); title_bar(s, "PART 7 · 結合架構", "結合架構：確定性 KG 區塊 + RAG 合成 = 超集")
# question
qn_box = box(s, Inches(0.7), Inches(1.6), Inches(2.4), Inches(0.9), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tf = qn_box.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text="使用者提問"; _set_font(r,15,True,WHITE)
# router
rt = box(s, Inches(3.7), Inches(1.6), Inches(2.4), Inches(0.9), fill=BLUE, shape=MSO_SHAPE.HEXAGON)
tf=rt.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text="混合路由\n判定題型"; _set_font(r,13,True,WHITE)
arrow(s, Inches(3.12), Inches(1.95), Inches(0.55), Inches(0.2), color=GRAY)
# two branches
box(s, Inches(6.7), Inches(0.9), Inches(6.0), Inches(1.5), fill=RGBColor(0xEC,0xFD,0xF5), line=GREEN)
text(s, Inches(6.9), Inches(1.0), Inches(5.6), Inches(0.4), [[("內容題 → 純 RAG", 15, True, GREEN)]])
text(s, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.85), [[("向量檢索 + 重排 + 小找大 + Grounded 合成 → 完整內容與數據", 13, False, DARK)]])
box(s, Inches(6.7), Inches(2.65), Inches(6.0), Inches(2.6), fill=RGBColor(0xF3,0xF0,0xFF), line=RGBColor(0x7A,0x3E,0xA8))
text(s, Inches(6.9), Inches(2.75), Inches(5.6), Inches(0.4), [[("關係題 → KG + RAG（超集）", 15, True, RGBColor(0x7A,0x3E,0xA8))]])
text(s, Inches(6.9), Inches(3.2), Inches(5.6), Inches(2.0),
     [[("① KG 確定性區塊（骨架）", 13.5, True, NAVY)],
      [("　版本家族 / 引用 / 被引用 / 取代鏈 — 直接讀圖，不漏邊", 12.5, False, DARK)],
      [("② ＋ RAG 合成（血肉）", 13.5, True, NAVY)],
      [("　補上文件內容、程序、背景說明", 12.5, False, DARK)],
      [("＝ Agent 答案 ⊇ 純 RAG", 13.5, True, RGBColor(0x7A,0x3E,0xA8))]])
arrow(s, Inches(6.12), Inches(1.55), Inches(0.55), Inches(0.2), color=GREEN)
arrow(s, Inches(6.12), Inches(3.85), Inches(0.55), Inches(0.2), color=RGBColor(0x7A,0x3E,0xA8))
box(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.2), fill=BG, line=LGRAY)
text(s, Inches(0.95), Inches(5.68), Inches(11.4), Inches(1.0),
     [[("核心理念：KG 給「結構化關係」確定且不漏，RAG 給「文件內容」完整且可溯源；兩者相加，", 14, False, DARK)],
      [("關係題的答案同時包含『精確引用清單（含章節出處）』＋『程序內容』，嚴格優於純 RAG。", 14, True, NAVY)]])

# real effect (relationship answer)
img_slide("PART 7 · 結合", "實際效果：關係查詢的「超集」答案",
          os.path.join(DECK, "11_rel_answer.png"),
          "問『510.7 引用哪些規範』→ KG 確定性區塊（6 個外部規範 + 章節出處）＋ RAG 補文件內容。",
          hl=[(0.40, 0.13, 0.56, 0.30, "KG 確定性引用清單（含章節）", "left"),
              (0.40, 0.45, 0.56, 0.20, "RAG 補充文件內容", "left")])

# ---- Part 8 AS-IS / TO-BE ----
section("08", "AS-IS vs TO-BE", "結合前 vs 結合後的回應效果")
s = slide(); title_bar(s, "PART 8 · 效果比較", "AS-IS（純 RAG） vs TO-BE（RAG × KG）")
headers = ["問題", "AS-IS：純 RAG", "TO-BE：RAG × KG"]
data = [
    ["510.7 引用哪些規範？", "模糊片段，漏掉完整 reference 清單", "6 個外部規範 + 章節出處 + 程序內容（超集）"],
    ["MIL-HDBK-310 ↔ MIL-STD-210 關係？", "「查無相關資料」", "引用 210B/210C；列出 47 處被引用"],
    ["MIL-STD-810 版本演進？", "只認得當前版", "完整版本家族 810→810B→…→810H"],
    ["810H 有哪些測試方法？", "亂列章節片段（SCOPE/定義…）", "完整 30 個 Method（500.6–528.1）"],
    ["MIL-DTL-901 被誰引用？", "含糊帶過", "明確：被 Method 516.8 §1.3/§4.6.9.3 引用"],
]
tl, tt, tw, th = Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.0)
tbl = s.shapes.add_table(len(data)+1, 3, tl, tt, tw, th).table
tbl.columns[0].width = Inches(3.4); tbl.columns[1].width = Inches(4.0); tbl.columns[2].width = Inches(4.85)
for j, hd in enumerate(headers):
    c = tbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    c.text = ""; p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = hd; _set_font(rr, 14, True, WHITE)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
for i, row in enumerate(data, start=1):
    for j, val in enumerate(row):
        c = tbl.cell(i, j)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF2,0xF5,0xF9)
        c.text = ""; tf = c.text_frame; tf.word_wrap = True; c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        if j == 0:
            rr = p.add_run(); rr.text = val; _set_font(rr, 12.5, True, NAVY)
        elif j == 1:
            rr = p.add_run(); rr.text = "✗ "; _set_font(rr, 12.5, True, RED)
            rr2 = p.add_run(); rr2.text = val; _set_font(rr2, 12.5, False, RGBColor(0x99,0x3D,0x2E))
        else:
            rr = p.add_run(); rr.text = "✓ "; _set_font(rr, 12.5, True, GREEN)
            rr2 = p.add_run(); rr2.text = val; _set_font(rr2, 12.5, False, RGBColor(0x1F,0x5C,0x2E))
text(s, Inches(0.55), Inches(6.75), Inches(12), Inches(0.5),
     [[("同一題、同一語料，差別只在「有沒有 KG」。", 14, True, NAVY)]], align=PP_ALIGN.CENTER)

# source + pdf + analysis
s = slide(); title_bar(s, "PART 8 · 效果比較", "TO-BE 加值：來源可溯 → PDF 預覽 → AI 頁面分析")
place_img(s, os.path.join(SHOTS, "src_pdf.png"), Inches(0.4), Inches(1.55), Inches(4.1), Inches(4.1))
place_img(s, os.path.join(SHOTS, "analyze_single.png"), Inches(4.65), Inches(1.55), Inches(4.1), Inches(4.1))
place_img(s, os.path.join(SHOTS, "analyze10_result.png"), Inches(8.9), Inches(1.55), Inches(4.0), Inches(4.1))
for cx, cap in [(Inches(0.4), "① 點來源 → 開啟 PDF 該頁"), (Inches(4.65), "② 單頁 AI 解析（VL 讀圖）"), (Inches(8.9), "③ 多頁分析（最多 10 頁、跨頁綜合）")]:
    text(s, cx, Inches(5.8), Inches(4.1), Inches(0.5), [[(cap, 13, True, NAVY)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
     [[("不只給答案，還能一路追到原始 PDF，並用視覺模型深入解析頁面圖表 — 查證閉環。", 14, False, GRAY)]], align=PP_ALIGN.CENTER)

# ---- Part 9 results ----
section("09", "成果與總結", "")
s = slide(); title_bar(s, "PART 9 · 成果", "量化成果")
stats = [("40 份", "軍規測試文件\n已匯入語料庫", BLUE), ("7,546", "文字 chunks\n(向量化)", GREEN),
         ("3,074 / 4,979", "KG 實體 / 關係", RGBColor(0x7A,0x3E,0xA8)), ("42 / 42", "回歸測試\n全數通過", ORANGE)]
x0 = Inches(0.7); y0 = Inches(1.8); cw = Inches(2.85); gap = Inches(0.25)
for i, (big, d, col) in enumerate(stats):
    cx = x0 + i*(cw+gap)
    box(s, cx, y0, cw, Inches(2.4), fill=WHITE, line=col, line_w=1.5)
    box(s, cx, y0, cw, Inches(0.16), fill=col)
    text(s, cx, y0+Inches(0.5), cw, Inches(0.8), [[(big, 30, True, col)]], align=PP_ALIGN.CENTER)
    text(s, cx, y0+Inches(1.45), cw, Inches(0.8), [[(ln, 14, False, DARK)] for ln in d.split("\n")], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.0), fill=BG, line=LGRAY)
text(s, Inches(0.95), Inches(4.78), Inches(11.4), Inches(1.7),
     [[("測試涵蓋 8 類問題（程序 / 數值 / 關係 / 列舉 / 版本 / 比較 / 應用 / 誠實），42 題全數通過。", 15, True, NAVY)],
      [("誠實題零幻覺：『810H 有無鋰電池安全測試』正確指向 MIL-STD-882；『EMI/EMC』正確指向 MIL-STD-461G（跨文件）。", 14, False, DARK)],
      [("混合路由 20/20 分類正確；關係題答案嚴格為純 RAG 的超集。", 14, False, DARK)]])

# takeaways
s = slide(); title_bar(s, "PART 9 · 總結", "Key Takeaways")
tk = [("RAG 調教解決『找得到、不亂編』", "重排 + 小找大 + 關閉思考 + Grounded 引用 + 誠實兜底。", BLUE),
      ("但跨文件關聯是 RAG 的天花板", "文字相似看不見『規範之間的關係』。", ORANGE),
      ("KG 解決『看得懂規範間的關係』", "引用網 + 版本鏈，確定性、可多跳、不漏項。", RGBColor(0x7A,0x3E,0xA8)),
      ("RAG × KG = 內容完整 + 關係精準", "KG 骨架 + RAG 血肉，混合路由各取所長 = 測試規範問答的最佳解。", GREEN)]
y = Inches(1.7)
for h, d, col in tk:
    box(s, Inches(0.7), y, Inches(11.9), Inches(1.1), fill=WHITE, line=LGRAY)
    box(s, Inches(0.7), y, Pt(7), Inches(1.1), fill=col)
    text(s, Inches(1.0), y+Inches(0.12), Inches(11.3), Inches(0.45), [[(h, 17, True, col)]])
    text(s, Inches(1.0), y+Inches(0.6), Inches(11.3), Inches(0.45), [[(d, 14, False, DARK)]])
    y += Inches(1.22)

# closing
s = slide(); box(s, 0, 0, SW, SH, fill=NAVY); box(s, Inches(0.9), Inches(3.9), Inches(1.2), Pt(4), fill=BLUE)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2),
     [[("RAG 找得到、KG 看得懂關係。", 34, True, WHITE)], [("兩者結合，才是測試規範智慧問答的最佳解。", 24, True, RGBColor(0xCF,0xE2,0xFF))]])
text(s, Inches(0.9), Inches(5.3), Inches(11), Inches(0.5), [[("Thank you　·　Q & A", 18, False, GRAY)]])

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
