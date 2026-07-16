# -*- coding: utf-8 -*-
"""RAG → KG 技術匯報 v2（繁體中文, 16:9, 給主管看得懂 + 有知識點）。

故事線：痛點 → 什麼是 RAG → 基本 RAG 的坑 → 如何一步步查得準（真實 MIL-STD-810H 案例）
       → 致命傷:跨文件關係 → 知識圖譜 KG → RAG×KG 混合 → AS-IS/TO-BE → 為何適合測試規範 → 未來。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

SHOTS = r"C:\Users\G635LXG\Downloads\RAG\AI_Document_V3\.shots"
DECK = os.path.join(SHOTS, "deck")
OUT = r"C:\Users\G635LXG\Downloads\RAG\RAG到知識圖譜_技術匯報.pptx"
FONT = "Microsoft JhengHei"

# palette
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x16, 0x77, 0xFF)
LBLUE = RGBColor(0xE7, 0xF0, 0xFF)
BG = RGBColor(0xF5, 0xF7, 0xFA)
RED = RGBColor(0xE0, 0x31, 0x31)
LRED = RGBColor(0xFC, 0xE9, 0xE9)
GREEN = RGBColor(0x2F, 0x9E, 0x44)
LGREEN = RGBColor(0xE6, 0xF5, 0xEA)
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
PURPLE = RGBColor(0x7A, 0x3E, 0xA8)
TEAL = RGBColor(0x0E, 0x7A, 0x8B)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xE5, 0xE7, 0xEB)
YELLOW = RGBColor(0xFF, 0xF4, 0xCC)
GOLD = RGBColor(0xB4, 0x7A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=DARK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=2):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for seg in para:
            txt, sz, bd, col = (seg + (None,) * 4)[:4]
            r = p.add_run(); r.text = txt
            _set_font(r, sz, bd or False, col or DARK)
    return tb


def title_bar(s, kicker, title):
    box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    box(s, 0, Inches(1.15), SW, Pt(3), fill=BLUE)
    if kicker:
        text(s, Inches(0.6), Inches(0.16), Inches(11.5), Inches(0.3), [[(kicker, 12, True, RGBColor(0x9E, 0xC5, 0xFF))]])
    text(s, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.6), [[(title, 26, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, items, l=Inches(0.7), t=Inches(1.5), w=Inches(12), h=Inches(5.6), sz=16, gap=8):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, txt, col, bold = (it + (0, "", DARK, False))[:4] if isinstance(it, tuple) else (0, it, DARK, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = lvl
        mark = "●  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark; _set_font(r, sz - 2, False, BLUE if lvl == 0 else GRAY)
        r2 = p.add_run(); r2.text = txt; _set_font(r2, sz if lvl == 0 else sz - 2, bold, col)
    return tb


def img_size(path):
    with Image.open(path) as im:
        return im.size


def place_img(s, path, bl, bt, bw, bh, border=True):
    iw, ih = img_size(path)
    r = min(bw / iw, bh / ih)
    w, h = int(iw * r), int(ih * r)
    l = bl + (bw - w) // 2; t = bt + (bh - h) // 2
    if border:
        box(s, l - Emu(9000), t - Emu(9000), w + Emu(18000), h + Emu(18000), fill=WHITE, line=LGRAY, line_w=1)
    s.shapes.add_picture(path, l, t, w, h)
    return l, t, w, h


def highlight(s, il, it, iw, ih, xf, yf, wf, hf, label="", lside="right"):
    l = il + int(iw * xf); t = it + int(ih * yf); w = int(iw * wf); h = int(ih * hf)
    hb = box(s, l, t, w, h, fill=None, line=RED, line_w=2.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        hb.adjustments[0] = 0.06
    except Exception:
        pass
    if label:
        tw, thh = Inches(2.2), Inches(0.34)
        if lside == "right":
            cl = l + w + Emu(45000); ct = t + h // 2 - thh // 2
        elif lside == "left":
            cl = max(Inches(0.05), l - tw - Emu(45000)); ct = t + h // 2 - thh // 2
        elif lside == "top":
            cl = l + w // 2 - tw // 2; ct = max(Inches(0.05), t - thh - Emu(30000))
        else:
            cl = l + w // 2 - tw // 2; ct = t + h + Emu(30000)
        tag = box(s, cl, ct, tw, thh, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = tag.text_frame; tf.word_wrap = True; tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = label; _set_font(rr, 11, True, WHITE)


def chip(s, l, t, w, h, label, fill, fg=WHITE, sz=12):
    c = box(s, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, sz, True, fg)
    return c


def arrow(s, l, t, w, h, color=BLUE):
    return box(s, l, t, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)


def marker(s, il, it, iw, ih, xf, yf, n, d=Inches(0.36)):
    cx = il + int(iw * xf) - d // 2; cy = it + int(ih * yf) - d // 2
    mk = box(s, cx, cy, d, d, fill=RED, line=WHITE, line_w=1.5, shape=MSO_SHAPE.OVAL)
    tf = mk.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); _set_font(r, 14, True, WHITE)


def pagenum(s, n):
    text(s, Inches(12.4), Inches(7.02), Inches(0.8), Inches(0.35), [[(str(n), 11, False, GRAY)]], align=PP_ALIGN.RIGHT)


PAGE = [0]


def newslide(kicker=None, title=None):
    s = slide()
    if title is not None:
        title_bar(s, kicker, title)
    PAGE[0] += 1
    pagenum(s, PAGE[0])
    return s


def section(num, title, sub=""):
    s = slide()
    box(s, 0, 0, SW, SH, fill=NAVY)
    box(s, Inches(0.9), Inches(3.05), Inches(1.2), Pt(5), fill=BLUE)
    text(s, Inches(0.9), Inches(2.05), Inches(3), Inches(1), [[(num, 60, True, RGBColor(0x2E, 0x5C, 0x9E))]])
    text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1), [[(title, 40, True, WHITE)]])
    if sub:
        text(s, Inches(0.95), Inches(4.55), Inches(11.5), Inches(0.6), [[(sub, 18, False, RGBColor(0xAF, 0xC6, 0xE0))]])
    PAGE[0] += 1
    pagenum(s, PAGE[0])
    return s


def case_slide(kicker, title, tech_tag, tech_desc, question, before, after, note=None):
    """強化前 → 強化後 一頁對照（真實規範問題）。"""
    s = newslide(kicker, title)
    # tech tag + one-line 原理
    chip(s, Inches(0.7), Inches(1.4), Inches(2.6), Inches(0.5), tech_tag, BLUE, sz=14)
    text(s, Inches(3.5), Inches(1.42), Inches(9.1), Inches(0.5), [[(tech_desc, 14.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    # question bar
    box(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(0.62), fill=NAVY)
    text(s, Inches(0.95), Inches(2.08), Inches(11.5), Inches(0.56),
         [[("工程師提問：", 13.5, True, RGBColor(0x9E, 0xC5, 0xFF)), (question, 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    # two panels
    py = Inches(2.95); ph = Inches(3.15); pw = Inches(5.86)
    # before (red)
    box(s, Inches(0.7), py, pw, ph, fill=LRED, line=RED, line_w=1.3)
    box(s, Inches(0.7), py, pw, Inches(0.6), fill=RED)
    text(s, Inches(0.9), py + Inches(0.1), pw - Inches(0.4), Inches(0.42), [[("✗  強化前（基本 RAG）", 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.95), py + Inches(0.78), pw - Inches(0.5), ph - Inches(0.95), [[(before, 14, False, DARK)]])
    # after (green)
    ax = Inches(0.7) + pw + Inches(0.21)
    box(s, ax, py, pw, ph, fill=LGREEN, line=GREEN, line_w=1.3)
    box(s, ax, py, pw, Inches(0.6), fill=GREEN)
    text(s, ax + Inches(0.2), py + Inches(0.1), pw - Inches(0.4), Inches(0.42), [[("✓  強化後（本系統）", 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, ax + Inches(0.25), py + Inches(0.78), pw - Inches(0.5), ph - Inches(0.95), [[(after, 14, False, DARK)]])
    # takeaway
    if note:
        box(s, Inches(0.7), Inches(6.32), Inches(11.93), Inches(0.62), fill=YELLOW, line=GOLD, line_w=1)
        text(s, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.56),
             [[("💡 知識點：", 13.5, True, GOLD), (note, 14, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    return s


# =====================================================================
# 1 · 封面
# =====================================================================
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(4.7), SW, Pt(4), fill=BLUE)
box(s, 0, 0, Inches(0.22), SH, fill=BLUE)
text(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.4), [[("測試規範智慧問答系統 · 技術匯報", 16, True, RGBColor(0x9E, 0xC5, 0xFF))]])
text(s, Inches(0.9), Inches(2.35), Inches(11.6), Inches(1.7),
     [[("從 RAG 到知識圖譜", 46, True, WHITE)], [("打造一個真正查得準、有關聯的規範知識庫", 28, True, RGBColor(0xCF, 0xE2, 0xFF))]])
text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.3),
     [[("內容查得準 × 關係查得到 × 答案可溯源", 18, False, RGBColor(0xCF, 0xDB, 0xEA))],
      [("以 MIL-STD-810H 軍規環境測試標準為例", 15, False, GRAY)]])
PAGE[0] += 1

# =====================================================================
# 2 · 痛點
# =====================================================================
s = newslide("為什麼要做這件事", "現況痛點：查一條測試規範，有多難？")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
     [[("一位測試工程師要回答「這個產品該做哪些環境測試、條件是什麼」，今天得靠人工翻閱大量 PDF：", 15.5, False, DARK)]])
pains = [
    ("📚", "文件又多又厚", "單一份 MIL-STD-810H 就上千頁，全套軍規/國際標準有數十份，關鍵條件藏在深處。", ORANGE),
    ("🔗", "交叉引用如迷宮", "一個測試方法動輒引用 5～10 份外部規範，要人工一份份追下去，極易遺漏。", RED),
    ("🕐", "版本新舊難分", "同一規範有 A/B/C 多版，哪版取代哪版、該用哪版，稍不留意就用錯。", PURPLE),
]
cw = Inches(3.95); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(2.15)
for i, (ic, h, b, col) in enumerate(pains):
    cx = x0 + i * (cw + gap)
    box(s, cx, y0, cw, Inches(3.5), fill=WHITE, line=LGRAY, line_w=1.2)
    box(s, cx, y0, cw, Pt(7), fill=col)
    text(s, cx, y0 + Inches(0.35), cw, Inches(0.9), [[(ic, 40, False, col)]], align=PP_ALIGN.CENTER)
    text(s, cx + Inches(0.25), y0 + Inches(1.35), cw - Inches(0.5), Inches(0.5), [[(h, 19, True, col)]], align=PP_ALIGN.CENTER)
    text(s, cx + Inches(0.3), y0 + Inches(1.95), cw - Inches(0.6), Inches(1.4), [[(b, 14, False, DARK)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(5.95), Inches(11.93), Inches(0.7), fill=NAVY)
text(s, Inches(0.95), Inches(5.99), Inches(11.5), Inches(0.62),
     [[("結果：查得慢、查不全、還可能查錯 — 而測試規範「答錯的代價」很高。", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# 3 · 目標願景
# =====================================================================
s = newslide("我們的目標", "把「翻幾百頁」變成「問一句話」")
text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.6),
     [[("目標：工程師用自然語言問一句，系統回覆有出處、可驗證、能一鍵點回 PDF 原頁的答案。", 16, False, DARK)]])
# before column
box(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(3.9), fill=BG, line=LGRAY)
box(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(0.7), fill=GRAY)
text(s, Inches(0.9), Inches(2.42), Inches(5.1), Inches(0.55), [[("今天：人工查閱", 18, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [(0, "翻上千頁 PDF、逐份追引用", DARK, False), (0, "靠經驗記憶，換人就斷層", DARK, False),
            (0, "跨規範關係全靠人腦拼湊", DARK, False), (0, "答案無出處，難查證、難稽核", DARK, False)],
        l=Inches(0.95), t=Inches(3.25), w=Inches(5.1), sz=15, gap=12)
arrow(s, Inches(6.35), Inches(4.0), Inches(0.65), Inches(0.55), color=BLUE)
# after column
box(s, Inches(7.15), Inches(2.35), Inches(5.48), Inches(3.9), fill=LBLUE, line=BLUE, line_w=1.3)
box(s, Inches(7.15), Inches(2.35), Inches(5.48), Inches(0.7), fill=BLUE)
text(s, Inches(7.35), Inches(2.42), Inches(5.1), Inches(0.55), [[("目標：問一句話", 18, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [(0, "自然語言提問，秒級回覆", DARK, False), (0, "答案每句附來源頁碼，可點開 PDF", DARK, False),
            (0, "跨規範引用/版本自動追出來", DARK, False), (0, "找不到就明說，絕不亂編", DARK, False)],
        l=Inches(7.4), t=Inches(3.25), w=Inches(5.1), sz=15, gap=12)

# =====================================================================
# Part 1 divider
# =====================================================================
section("PART 1", "什麼是 RAG？", "先用三頁，讓大家對「檢索增強生成」有共識")

# 4 · LLM 會亂編
s = newslide("PART 1 · 什麼是 RAG", "問題根源：大型語言模型（LLM）會「一本正經地亂編」")
text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.6),
     [[("LLM（像 ChatGPT）很會講話，但它只憑「訓練時讀過的網路資料」作答 —", 16, False, DARK)],
      [("它沒讀過我們公司的規範文件，遇到不知道的，會用流暢的語氣「編一個」。這叫幻覺（Hallucination）。", 16, False, DARK)]])
# analogy
box(s, Inches(0.7), Inches(2.85), Inches(11.93), Inches(1.5), fill=YELLOW, line=GOLD, line_w=1)
text(s, Inches(0.95), Inches(2.98), Inches(11.5), Inches(1.3),
     [[("🧠 打個比方：", 16, True, GOLD), ("LLM 就像一位口才極好、見多識廣的顧問 —", 16, True, DARK)],
      [("     但他從沒看過「你們公司的內部規範」。你問他細節，他不會說「我不知道」，", 15.5, False, DARK)],
      [("     而是憑印象講一套聽起來很專業、其實可能是錯的答案。", 15.5, False, DARK)]])
# consequences
box(s, Inches(0.7), Inches(4.65), Inches(11.93), Inches(1.9), fill=LRED, line=RED, line_w=1.2)
text(s, Inches(0.95), Inches(4.78), Inches(11.5), Inches(0.5), [[("⚠️ 直接把 LLM 用在測試規範查詢，會踩到三個致命問題：", 15.5, True, RED)]])
bullets(s, [(0, "編造條件：把 507.6 濕度測試的溫濕度條件講錯，工程師照做就出事", DARK, False),
            (0, "看不到內部文件：模型內建知識裡根本沒有我們的私有規範全文", DARK, False),
            (0, "無法查證：講得頭頭是道，卻給不出「這句話出自哪一份、哪一頁」", DARK, False)],
        l=Inches(1.0), t=Inches(5.28), w=Inches(11.3), sz=14.5, gap=6)

# 5 · RAG = 開卷考試
s = newslide("PART 1 · 什麼是 RAG", "RAG = 讓 AI「先查資料、再作答」的開卷考試")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.55),
     [[("RAG（檢索增強生成）：回答前，先到我們的文件庫「檢索」出相關片段，再讓 LLM「只根據這些片段」作答並附出處。", 15, False, DARK)]])
steps = ["文件\nPDF / 規範", "切塊\n拆成小段", "向量化\n轉成語意座標", "向量庫\nFAISS", "提問檢索\n找最相關 Top-K", "生成答案\n附來源頁碼"]
n = len(steps); x = Inches(0.55); y = Inches(2.45); bw = Inches(1.78); bh = Inches(1.35); g = Inches(0.28)
cols = [NAVY, BLUE, BLUE, TEAL, ORANGE, PURPLE]
for i, st in enumerate(steps):
    bx = x + i * (bw + g)
    c = box(s, bx, y, bw, bh, fill=cols[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for j, ln in enumerate(st.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, 14 if j == 0 else 10.5, j == 0, WHITE)
    if i < n - 1:
        arrow(s, bx + bw + Emu(15000), y + bh // 2 - Inches(0.12), g - Emu(30000), Inches(0.24), color=GRAY)
box(s, Inches(0.7), Inches(4.15), Inches(11.93), Inches(0.75), fill=LBLUE, line=BLUE, line_w=1)
text(s, Inches(0.95), Inches(4.2), Inches(11.5), Inches(0.65),
     [[("關鍵差別：答案來自「真的檢索到的文件片段」，每一句都能標上來源頁碼 → 可溯源、可查證、不亂編。", 15, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.7), Inches(5.1), Inches(11.93), Inches(1.4), fill=BG, line=LGRAY)
text(s, Inches(0.95), Inches(5.24), Inches(11.5), Inches(1.15),
     [[("🔧 本系統實際的 RAG 管線：", 14.5, True, NAVY)],
      [("pdfminer 抽文字  →  切塊  →  qwen3-embedding 向量化  →  FAISS 向量庫  →  向量檢索 + 交叉編碼器重排  →  gemma 生成帶引用答案", 13.5, False, DARK)],
      [("（下一部分會說明,光有這條基本管線還「查不夠準」,我們做了哪些強化。）", 13, False, GRAY)]])

# 6 · 為何需要 RAG
s = newslide("PART 1 · 什麼是 RAG", "為什麼「測試規範查詢」特別需要 RAG")
cards = [("① 消除幻覺", "答案必須有檢索到的文件撐腰，LLM 不能再憑印象亂編條件。", BLUE),
         ("② 讀得到私有文件", "軍規/測試標準無法內建於通用模型，RAG 讓系統即時讀取我們自己的語料。", GREEN),
         ("③ 可溯源可稽核", "每個答案附來源頁碼，點開就是 PDF 原文 — 這是合規與品保場景的硬需求。", ORANGE)]
cw = Inches(3.95); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(1.7)
for i, (h, b, col) in enumerate(cards):
    cx = x0 + i * (cw + gap)
    box(s, cx, y0, cw, Inches(3.55), fill=WHITE, line=LGRAY, line_w=1.2)
    box(s, cx, y0, cw, Inches(0.85), fill=col)
    text(s, cx + Inches(0.25), y0 + Inches(0.18), cw - Inches(0.5), Inches(0.5), [[(h, 19, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + Inches(0.3), y0 + Inches(1.15), cw - Inches(0.6), Inches(2.2), [[(b, 15.5, False, DARK)]])
box(s, Inches(0.7), Inches(5.65), Inches(11.93), Inches(0.85), fill=NAVY)
text(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(0.75),
     [[("→ 對「精確、可查證、來自大量私有 PDF」的測試規範場景，RAG 是最基本的門檻。但——它只是起點。", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)

# 7 · 基本 RAG 的四大坑
s = newslide("PART 1 · 什麼是 RAG", "但基本 RAG 會踩的四大坑：「查得到」≠「查得準」")
pits = [
    ("坑 1｜切太爛", "文件被機械式切段，一句完整條件被切成兩半，檢索到半句 → 答案殘缺。", ORANGE),
    ("坑 2｜找錯段", "只比對字面關鍵字：問「濕度」卻對不上文件寫的「humidity / moisture」→ 漏檢。", RED),
    ("坑 3｜沒出處", "給了答案卻標不出處，工程師無法查證，等於還是不敢用。", PURPLE),
    ("坑 4｜關係查不到", "「510.7 引用哪些規範?」「哪版取代哪版?」— 關係散在跨文件，基本 RAG 完全無感。", NAVY),
]
x0 = Inches(0.7); y0 = Inches(1.65); cw = Inches(5.86); ch = Inches(1.95); gx = Inches(0.21); gy = Inches(0.22)
for i, (h, b, col) in enumerate(pits):
    r, c = divmod(i, 2)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    box(s, cx, cy, cw, ch, fill=WHITE, line=LGRAY, line_w=1.2)
    box(s, cx, cy, Pt(7), ch, fill=col)
    text(s, cx + Inches(0.3), cy + Inches(0.22), cw - Inches(0.5), Inches(0.5), [[(h, 18, True, col)]])
    text(s, cx + Inches(0.35), cy + Inches(0.85), cw - Inches(0.6), Inches(1.0), [[(b, 14.5, False, DARK)]])
text(s, Inches(0.7), Inches(6.35), Inches(11.93), Inches(0.6),
     [[("坑 1～3 靠「把 RAG 做準」解決（Part 2）；坑 4 是致命傷,要靠「知識圖譜」解決（Part 3）。", 14.5, True, RED)]], align=PP_ALIGN.CENTER)

# =====================================================================
# Part 2 divider
# =====================================================================
section("PART 2", "如何一步步把 RAG 做準", "以真實 MIL-STD-810H 問題，看「強化前 → 強化後」的差別")

# 8 · 五大強化總覽
s = newslide("PART 2 · 把 RAG 做準", "五道強化：把「查得到」磨成「查得準」")
enh = ["① 聰明切塊\nSmall-to-Big", "② 語意檢索\nEmbedding", "③ 重排序\nRerank", "④ 強制出處\nGrounded", "⑤ 防幻覺\n找不到就說"]
n = len(enh); x = Inches(0.6); y = Inches(2.0); bw = Inches(2.25); bh = Inches(1.5); g = Inches(0.18)
cols = [NAVY, BLUE, TEAL, ORANGE, GREEN]
for i, st in enumerate(enh):
    bx = x + i * (bw + g)
    c = box(s, bx, y, bw, bh, fill=cols[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for j, ln in enumerate(st.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, 15 if j == 0 else 11, True, WHITE)
    if i < n - 1:
        arrow(s, bx + bw + Emu(4000), y + bh // 2 - Inches(0.1), g - Emu(8000), Inches(0.2), color=GRAY)
box(s, Inches(0.7), Inches(3.9), Inches(11.93), Inches(0.7), fill=LBLUE, line=BLUE, line_w=1)
text(s, Inches(0.95), Inches(3.95), Inches(11.5), Inches(0.6),
     [[("每一道強化都對應前面的一個坑。接下來用真實的 810H 問題,逐一展示「強化前 vs 強化後」。", 15, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
# map to pits
mp = [("① 聰明切塊", "解坑 1：不再切斷條件", ORANGE), ("② 語意檢索", "解坑 2：懂同義、不漏檢", RED),
      ("③ 重排序", "解坑 2：把最相關的推到最前", RED), ("④ 強制出處", "解坑 3：每句可點回 PDF", PURPLE),
      ("⑤ 防幻覺", "守底線：沒有就說沒有", NAVY)]
x0 = Inches(0.7); cw2 = Inches(2.3); y2 = Inches(4.85)
for i, (h, d, col) in enumerate(mp):
    cx = x0 + i * (cw2 + Inches(0.11))
    box(s, cx, y2, cw2, Inches(1.4), fill=WHITE, line=LGRAY)
    box(s, cx, y2, cw2, Pt(6), fill=col)
    text(s, cx + Inches(0.12), y2 + Inches(0.2), cw2 - Inches(0.24), Inches(0.4), [[(h, 14, True, col)]], align=PP_ALIGN.CENTER)
    text(s, cx + Inches(0.14), y2 + Inches(0.72), cw2 - Inches(0.28), Inches(0.6), [[(d, 12.5, False, DARK)]], align=PP_ALIGN.CENTER)

# 9 · 強化1 聰明切塊
case_slide(
    "PART 2 · 強化 ①", "聰明切塊：先精準命中，再自動補回上下文",
    "Small-to-Big", "命中小片段後，自動把它所在的完整段落一起帶給 LLM。",
    "溫度衝擊測試（Method 503.7）的高低溫停留條件是什麼？",
    "文件被固定長度切段，「停留至溫度穩定後維持…」這句被切成兩半，只檢索到前半。\n\nLLM 只拿到殘缺片段 → 答案漏掉「維持時間 / 循環次數」等關鍵條件，工程師照做會出錯。",
    "先用小片段精準命中「停留條件」那一句,再自動往外補回它所屬的完整程序段落。\n\nLLM 拿到的是完整條件 → 停留、穩定判定、循環次數一次講齊,不再殘缺。",
    "切太細會斷句、切太粗會混雜；小塊命中、大塊回答,兩者兼得。",
)

# 10 · 強化2 語意檢索
case_slide(
    "PART 2 · 強化 ②", "語意檢索：找的是「意思」，不是「關鍵字」",
    "Embedding 向量", "把文字轉成語意座標,意思相近的靠在一起,不靠字面比對。",
    "濕度測試怎麼進行？",
    "傳統關鍵字檢索：使用者打「濕度」,但文件原文寫的是「humidity / moisture / warm-humid」。\n\n字面對不上 → 直接漏檢,回「查無相關內容」,或抓到不相干段落。",
    "語意向量理解「濕度 ≈ humidity ≈ moisture ≈ 潮濕環境」,不管用哪種說法都能命中。\n\n正確定位到 Method 507.6 濕度測試,回出真正的程序內容。",
    "工程師的口語問法,和規範的英文術語常對不上;語意檢索跨過了這道語言鴻溝。",
)

# 11 · 強化3 重排序
case_slide(
    "PART 2 · 強化 ③", "重排序：把「最該看的那段」推到第一位",
    "Cross-Encoder Rerank", "對初步檢索結果做第二次精算,依「與問題的真正相關度」重新排序。",
    "濕度測試的實際操作步驟有哪些？",
    "向量檢索一次撈回 Top-5,裡面混了「濕度的定義」「適用範圍」「引用清單」等相關但非程序的段落。\n\nLLM 分不清輕重,可能拿「定義」當「步驟」回答 → 答非所問。",
    "交叉編碼器逐段重算相關度,把「507.6 程序步驟」那段推到第 1 名,雜訊段落往後排。\n\nLLM 先看到最對的內容 → 步驟講得精準、不被雜訊帶偏。",
    "檢索負責「大致找到」,重排負責「精準排序」;兩段式篩選,準確度明顯提升。",
)

# 12 · 強化4 強制出處
s = newslide("PART 2 · 強化 ④", "強制出處：每一句答案，都能一鍵點回 PDF 原頁")
chip(s, Inches(0.7), Inches(1.4), Inches(2.6), Inches(0.5), "Grounded Synthesis", BLUE, sz=13)
text(s, Inches(3.5), Inches(1.42), Inches(9.1), Inches(0.5), [[("要求 LLM「只能根據檢索到的片段作答」,並在每個重點標註來源。", 14.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.7), Inches(2.05), Inches(5.9), Inches(4.3), fill=BG, line=LGRAY)
text(s, Inches(0.9), Inches(2.2), Inches(5.5), Inches(0.5), [[("實際查詢後,答案下方就是可點的來源清單：", 14.5, True, NAVY)]])
bullets(s, [(0, "每個重點對應一份文件 + 頁碼", DARK, False),
            (0, "點一下 → 直接開啟該 PDF 的那一頁", DARK, False),
            (0, "工程師可立即核對原文,不必全盤相信 AI", DARK, False),
            (0, "稽核 / 品保時,答案來源一目了然", DARK, False)],
        l=Inches(0.95), t=Inches(2.85), w=Inches(5.4), sz=14.5, gap=12)
box(s, Inches(0.9), Inches(5.35), Inches(5.5), Inches(0.85), fill=YELLOW, line=GOLD, line_w=1)
text(s, Inches(1.05), Inches(5.42), Inches(5.2), Inches(0.72),
     [[("💡 這一步把 RAG 從「聽起來對」升級成「可被驗證」,是合規場景能落地的關鍵。", 13.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
il, it, iw, ih = place_img(s, os.path.join(DECK, "02_qa_main.png"), Inches(6.85), Inches(2.05), Inches(5.75), Inches(4.3))
highlight(s, il, it, iw, ih, 0.60, 0.12, 0.38, 0.30, "答案 + 可點來源", lside="bottom")

# 13 · 強化5 防幻覺
case_slide(
    "PART 2 · 強化 ⑤", "防幻覺：找不到，就老實說找不到",
    "拒答護欄", "檢索不到足夠依據時,系統寧可說「知識庫沒有」,也不編造。",
    "這批貨要符合 ASTM B117 鹽霧試驗,條件是什麼？（註：ASTM 為付費規範,未收錄）",
    "基本 RAG 沒有護欄,檢索不到也硬答:憑印象生出一段看似專業、實則可能錯誤的 B117 條件。\n\n工程師若照單全收,風險極高 — 這正是 LLM 幻覺最危險的地方。",
    "系統偵測到知識庫沒有 ASTM B117 全文,明確回覆:「未收錄該規範」,並主動指向館內可用的\nMIL-STD-810H Method 509 鹽霧對應內容作為替代參考。",
    "會說「我不知道」的 AI,才是可信任的 AI;一次亂編,信任就崩了。",
)

# 14 · 小結
s = newslide("PART 2 · 小結", "五道強化 → 單篇規範,終於「查得準」")
text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
     [[("經過五道強化,對「單一份規範的內容問題」,系統已能穩定給出完整、正確、可溯源的答案。", 16, True, NAVY)]])
rows = [("坑 1 切太爛", "① 聰明切塊", "條件不再被切斷"),
        ("坑 2 找錯段", "② 語意檢索 + ③ 重排序", "懂同義、排序準"),
        ("坑 3 沒出處", "④ 強制出處", "每句可點回 PDF"),
        ("— 幻覺風險", "⑤ 防幻覺護欄", "沒有就說沒有")]
y = Inches(2.4); rh = Inches(0.82); x0 = Inches(0.9)
w1, w2, w3 = Inches(3.3), Inches(4.6), Inches(3.6)
hdr = [("原本的坑", w1, RED), ("對應的強化", w2, BLUE), ("結果", w3, GREEN)]
cx = x0
for h, w, col in hdr:
    box(s, cx, y, w, Inches(0.55), fill=col)
    text(s, cx, y, w, Inches(0.55), [[(h, 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i, (a, b, c) in enumerate(rows):
    ry = y + Inches(0.55) + i * rh
    fill = WHITE if i % 2 == 0 else BG
    cx = x0
    for val, w in ((a, w1), (b, w2), (c, w3)):
        box(s, cx, ry, w, rh, fill=fill, line=LGRAY)
        bold = w is w2
        text(s, cx + Inches(0.15), ry, w - Inches(0.3), rh, [[(val, 14.5, bold, DARK)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        cx += w
box(s, Inches(0.7), Inches(6.35), Inches(11.93), Inches(0.7), fill=RED)
text(s, Inches(0.95), Inches(6.39), Inches(11.5), Inches(0.62),
     [[("但——就算每一份都查得再準,還有一種問題,單篇 RAG 永遠答不出來 →", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# Part 3 divider
# =====================================================================
section("PART 3", "從「內容」到「關係」：知識圖譜", "RAG 的致命傷,與 KG 如何補上這一塊")

# 15 · 致命傷
s = newslide("PART 3 · 致命傷", "純 RAG 的天花板：關係題，它答不出來")
box(s, Inches(0.7), Inches(1.45), Inches(11.93), Inches(0.7), fill=NAVY)
text(s, Inches(0.95), Inches(1.49), Inches(11.5), Inches(0.62),
     [[("工程師提問：", 13.5, True, RGBColor(0x9E, 0xC5, 0xFF)), ("沙塵測試 Method 510.7 到底引用了哪些外部規範？", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
# before
box(s, Inches(0.7), Inches(2.4), Inches(5.86), Inches(3.4), fill=LRED, line=RED, line_w=1.3)
box(s, Inches(0.7), Inches(2.4), Inches(5.86), Inches(0.6), fill=RED)
text(s, Inches(0.9), Inches(2.5), Inches(5.5), Inches(0.42), [[("✗  純 RAG 只會「讀一段文字」", 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.95), Inches(3.15), Inches(5.5), Inches(2.55),
     [[("它檢索到 510.7 的某一段,那段剛好只提到 2 份規範,就只答 2 份。", 14, False, DARK)],
      [("", 6, False, DARK)],
      [("但「引用」這件事散落在整份方法的許多段落、甚至跨到別的文件 → 答案殘缺、時多時少,還可能憑印象補幾個(幻覺)。", 14, False, DARK)]])
# after
box(s, Inches(6.77), Inches(2.4), Inches(5.86), Inches(3.4), fill=LGREEN, line=GREEN, line_w=1.3)
box(s, Inches(6.77), Inches(2.4), Inches(5.86), Inches(0.6), fill=GREEN)
text(s, Inches(6.97), Inches(2.5), Inches(5.5), Inches(0.42), [[("✓  知識圖譜「查關係」,給完整清單", 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(7.02), Inches(3.1), Inches(5.4), Inches(0.4), [[("系統一次列出 510.7 的完整引用（實測結果）：", 13.5, True, DARK)]])
refs = ["ASTM D185-07", "IEC 60068-2-68", "MIL-HDBK-310", "MIL-STD-210B", "MIL-STD-3033"]
for i, rf in enumerate(refs):
    rr, cc = divmod(i, 2)
    chip(s, Inches(7.02) + cc * Inches(2.75), Inches(3.55) + rr * Inches(0.55), Inches(2.6), Inches(0.44), rf, NAVY, sz=13)
text(s, Inches(7.02), Inches(5.25), Inches(5.4), Inches(0.5), [[("五份,一份不漏、不多不少、可驗證。", 14, True, GREEN)]])
text(s, Inches(0.7), Inches(6.0), Inches(11.93), Inches(0.9),
     [[("為什麼差這麼多?因為問的不是「內容」,而是「關係」— 而關係,不在任何單一段文字裡。", 16, True, RED)]], align=PP_ALIGN.CENTER)

# 16 · 為什麼
s = newslide("PART 3 · 致命傷", "為什麼關係題這麼難：關係散在「文件與文件之間」")
text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.6),
     [[("RAG 的世界觀是「一堆互不相干的文字片段」。它很會找片段,卻看不見片段之間的線。", 16, False, DARK)]])
box(s, Inches(0.7), Inches(2.3), Inches(5.86), Inches(3.7), fill=LRED, line=RED, line_w=1.2)
text(s, Inches(0.95), Inches(2.45), Inches(5.4), Inches(0.5), [[("RAG 眼中的文件庫", 17, True, RED)]])
bullets(s, [(0, "一段段孤立的文字", DARK, False),
            (0, "只能問:「哪段文字最像我的問題?」", DARK, False),
            (0, "看不到 A 規範「引用」B 規範", DARK, False),
            (0, "看不到 C 版「取代」B 版", DARK, False),
            (0, "關係題 → 只能瞎猜或漏答", RED, True)],
        l=Inches(1.0), t=Inches(3.05), w=Inches(5.3), sz=14.5, gap=11)
box(s, Inches(6.77), Inches(2.3), Inches(5.86), Inches(3.7), fill=LGREEN, line=GREEN, line_w=1.2)
text(s, Inches(7.02), Inches(2.45), Inches(5.4), Inches(0.5), [[("我們需要的世界觀", 17, True, GREEN)]])
bullets(s, [(0, "一張「規範關係網」", DARK, False),
            (0, "節點 = 規範 / 方法 / 章節", DARK, False),
            (0, "連線 = 引用 / 取代 / 要求 / 定義", DARK, False),
            (0, "問「510.7 引用誰」= 沿著線走一步", DARK, False),
            (0, "關係題 → 沿圖查,完整又確定", GREEN, True)],
        l=Inches(7.07), t=Inches(3.05), w=Inches(5.3), sz=14.5, gap=11)
text(s, Inches(0.7), Inches(6.2), Inches(11.93), Inches(0.6),
     [[("這張「關係網」,就是知識圖譜（Knowledge Graph, KG）。", 17, True, NAVY)]], align=PP_ALIGN.CENTER)

# 17 · KG 是什麼
s = newslide("PART 3 · 知識圖譜", "知識圖譜（KG）是什麼：把文件變成一張「關係網」")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.55),
     [[("KG = 從文件中抽出「實體」(節點) 和它們之間的「關係」(連線),組成一張可查詢的網。", 15.5, False, DARK)]])
# mini concept graph
gx, gy = Inches(1.3), Inches(2.5)
def node(cx, cy, label, col, w=Inches(2.0), h=Inches(0.7)):
    b = box(s, cx, cy, w, h, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, 12.5, True, WHITE)
    return cx + w // 2, cy + h // 2
def edge(x1, y1, x2, y2, label, col=GRAY):
    ln = s.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = col; ln.line.width = Pt(1.75); ln.shadow.inherit = False
    mx, my = (x1 + x2) // 2, (y1 + y2) // 2
    tg = box(s, mx - Inches(0.55), my - Inches(0.17), Inches(1.1), Inches(0.34), fill=WHITE, line=col, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tg.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, 10.5, True, col)
c1 = node(Inches(0.95), Inches(3.4), "MIL-STD-810H", NAVY)
c2 = node(Inches(4.5), Inches(2.35), "Method 510.7\n沙塵", ORANGE, h=Inches(0.8))
c3 = node(Inches(4.5), Inches(4.35), "Method 507.6\n濕度", TEAL, h=Inches(0.8))
c4 = node(Inches(8.2), Inches(1.7), "IEC 60068-2-68", BLUE)
c5 = node(Inches(8.2), Inches(3.0), "MIL-HDBK-310", GREEN)
c6 = node(Inches(8.2), Inches(4.55), "MIL-STD-210B", PURPLE)
c7 = node(Inches(10.9), Inches(2.35), "IEC 60068-2-52", RED)
edge(*c1, *c2, "包含")
edge(*c1, *c3, "包含")
edge(*c2, *c4, "引用", BLUE)
edge(*c2, *c5, "引用", BLUE)
edge(*c2, *c6, "引用", BLUE)
edge(*c4, *c7, "取代", RED)
text(s, Inches(0.7), Inches(5.75), Inches(11.93), Inches(1.1),
     [[("讀法:「810H 包含 510.7,510.7 引用 IEC 60068-2-68,而 IEC 60068-2-68 取代了 60068-2-52。」", 15, True, NAVY)],
      [("每一句關係都是圖上的一條線 → 關係題不再靠讀文字猜,而是沿著線走。", 14.5, False, DARK)]])

# 18 · 真實 KG 圖
s = newslide("PART 3 · 知識圖譜", "這不是概念圖：本系統的規範關係圖是真的")
il, it, iw, ih = place_img(s, os.path.join(DECK, "kg_zoom_310.png"), Inches(0.5), Inches(1.5), Inches(8.4), Inches(5.1))
highlight(s, il, it, iw, ih, 0.115, 0.10, 0.14, 0.22, "全庫統計", lside="bottom")
# stats panel
lx = Inches(9.15); ly = Inches(1.6)
box(s, lx, ly, Inches(3.75), Inches(5.0), fill=BG, line=LGRAY)
text(s, lx + Inches(0.25), ly + Inches(0.2), Inches(3.3), Inches(0.5), [[("目前圖譜規模（實測）", 15, True, NAVY)]])
stat = [("3,074", "個實體節點", BLUE), ("4,979", "條關係連線", GREEN), ("40", "份軍規文件已入圖", ORANGE)]
for i, (num, lab, col) in enumerate(stat):
    yy = ly + Inches(0.75) + i * Inches(0.95)
    text(s, lx + Inches(0.25), yy, Inches(3.3), Inches(0.5), [[(num, 30, True, col), ("  " + lab, 13, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, lx + Inches(0.2), ly + Inches(3.7), Inches(3.35), Inches(1.1), fill=WHITE, line=LGRAY)
text(s, lx + Inches(0.35), ly + Inches(3.8), Inches(3.05), Inches(0.95),
     [[("關係類型:引用 2,014、要求 351、", 12.5, False, DARK)], [("取代 43、衍生自 37、定義 3…", 12.5, False, DARK)],
      [("全部由系統自動從 PDF 抽取。", 12.5, True, NAVY)]])
text(s, Inches(0.5), Inches(6.68), Inches(8.4), Inches(0.5),
     [[("在「知識圖譜」頁輸入任一規範,即可展開它的引用網與版本鏈。", 13, False, GRAY)]], align=PP_ALIGN.CENTER)

# 19 · RAG 查內容 vs KG 查關係
s = newslide("PART 3 · 知識圖譜", "一句話記住：RAG 查「內容」，KG 查「關係」")
box(s, Inches(0.7), Inches(1.7), Inches(5.86), Inches(4.3), fill=LBLUE, line=BLUE, line_w=1.3)
box(s, Inches(0.7), Inches(1.7), Inches(5.86), Inches(0.8), fill=BLUE)
text(s, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.6), [[("RAG　查「內容」", 20, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [(0, "問題像:「濕度測試怎麼做?」", DARK, False),
            (0, "「503.7 的停留條件是什麼?」", DARK, False),
            (0, "答案在「某一段文字」裡", DARK, False),
            (0, "強項:讀懂內文、給細節、附出處", BLUE, True)],
        l=Inches(1.0), t=Inches(2.85), w=Inches(5.3), sz=15.5, gap=14)
box(s, Inches(6.77), Inches(1.7), Inches(5.86), Inches(4.3), fill=LGREEN, line=GREEN, line_w=1.3)
box(s, Inches(6.77), Inches(1.7), Inches(5.86), Inches(0.8), fill=GREEN)
text(s, Inches(6.97), Inches(1.8), Inches(5.5), Inches(0.6), [[("KG　查「關係」", 20, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [(0, "問題像:「510.7 引用哪些規範?」", DARK, False),
            (0, "「哪一版取代了舊版?」", DARK, False),
            (0, "答案在「文件與文件的連線」裡", DARK, False),
            (0, "強項:跨文件、列清單、追版本鏈", GREEN, True)],
        l=Inches(7.07), t=Inches(2.85), w=Inches(5.3), sz=15.5, gap=14)
box(s, Inches(0.7), Inches(6.2), Inches(11.93), Inches(0.7), fill=NAVY)
text(s, Inches(0.95), Inches(6.24), Inches(11.5), Inches(0.62),
     [[("兩者互補,缺一不可。下一步:讓系統自動判斷該用哪一個 →", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# 20 · 混合路由
# =====================================================================
s = newslide("PART 3 · RAG × KG 結合", "混合查詢:系統自動判斷「這題該走哪條路」")
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.55),
     [[("使用者不必懂技術。系統看問題的語氣,自動路由:內容題交給 RAG,關係題交給 KG(Agent)。", 15, False, DARK)]])
# router flow
box(s, Inches(5.0), Inches(2.25), Inches(3.3), Inches(0.85), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(5.0), Inches(2.25), Inches(3.3), Inches(0.85), [[("使用者問一句話", 16, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(4.85), Inches(3.35), Inches(3.6), Inches(0.75), fill=ORANGE, shape=MSO_SHAPE.DIAMOND)
text(s, Inches(4.85), Inches(3.35), Inches(3.6), Inches(0.75), [[("自動路由判斷", 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# left RAG
box(s, Inches(1.1), Inches(4.6), Inches(4.7), Inches(1.7), fill=LBLUE, line=BLUE, line_w=1.3)
text(s, Inches(1.3), Inches(4.72), Inches(4.3), Inches(0.5), [[("內容題 → 走 RAG", 16, True, BLUE)]])
text(s, Inches(1.3), Inches(5.25), Inches(4.3), Inches(0.95),
     [[("「怎麼做 / 條件是什麼 / 數值多少」", 13.5, False, DARK)], [("快、完整、附 PDF 出處", 13.5, True, DARK)]])
# right KG
box(s, Inches(7.55), Inches(4.6), Inches(4.7), Inches(1.7), fill=LGREEN, line=GREEN, line_w=1.3)
text(s, Inches(7.75), Inches(4.72), Inches(4.3), Inches(0.5), [[("關係題 → 走 KG / Agent", 16, True, GREEN)]])
text(s, Inches(7.75), Inches(5.25), Inches(4.3), Inches(0.95),
     [[("「引用誰 / 哪版取代 / 有哪些子項」", 13.5, False, DARK)], [("跨規範追查,給確定清單", 13.5, True, DARK)]])
ln1 = s.shapes.add_connector(2, Inches(5.3), Inches(4.1), Inches(3.45), Inches(4.6)); ln1.line.color.rgb = BLUE; ln1.line.width = Pt(2); ln1.shadow.inherit = False
ln2 = s.shapes.add_connector(2, Inches(8.0), Inches(4.1), Inches(9.9), Inches(4.6)); ln2.line.color.rgb = GREEN; ln2.line.width = Pt(2); ln2.shadow.inherit = False
text(s, Inches(0.7), Inches(6.55), Inches(11.93), Inches(0.5),
     [[("預設走「混合」,每次查詢都會標示實際走了哪一邊 → 兼顧速度與完整性,使用者零學習成本。", 14, True, NAVY)]], align=PP_ALIGN.CENTER)

# 21 · 三模式主介面
s = newslide("PART 3 · RAG × KG 結合", "主介面:三種模式一鍵切換")
il, it, iw, ih = place_img(s, os.path.join(DECK, "02_qa_main.png"), Inches(0.45), Inches(1.5), Inches(7.7), Inches(5.0))
feat = [(0.345, 0.088, "模式切換", "純RAG / 混合(預設) / Agent"),
        (0.23, 0.305, "自然語言提問", "打一句話,Ctrl+Enter 送出"),
        (0.23, 0.45, "範圍篩選", "可限定分類 / 專案 / 文件"),
        (0.23, 0.66, "來源筆數", "控制檢索回傳段落數"),
        (0.62, 0.12, "結果區", "答案 + 可點來源 + 追問")]
for i, (xf, yf, _, _) in enumerate(feat, 1):
    marker(s, il, it, iw, ih, xf, yf, i)
lx = Inches(8.4); ly = Inches(1.6)
box(s, lx, ly, Inches(4.5), Inches(4.9), fill=BG, line=LGRAY)
tb = s.shapes.add_textbox(lx + Inches(0.2), ly + Inches(0.2), Inches(4.1), Inches(4.5)); tf = tb.text_frame; tf.word_wrap = True
for i, (_, _, h, d) in enumerate(feat, 1):
    p = tf.paragraphs[0] if i == 1 else tf.add_paragraph(); p.space_after = Pt(4)
    r = p.add_run(); r.text = f"{i}  "; _set_font(r, 15, True, RED)
    r2 = p.add_run(); r2.text = h; _set_font(r2, 15, True, NAVY)
    p2 = tf.add_paragraph(); p2.space_after = Pt(12)
    r3 = p2.add_run(); r3.text = "      " + d; _set_font(r3, 12.5, False, DARK)

# 22 · 關係題確定性答案（真實截圖）
s = newslide("PART 3 · RAG × KG 結合", "關係題實測:Agent 沿圖查,給出確定清單")
il, it, iw, ih = place_img(s, os.path.join(DECK, "11_rel_answer.png"), Inches(0.5), Inches(1.5), Inches(8.7), Inches(5.05))
highlight(s, il, it, iw, ih, 0.485, 0.40, 0.30, 0.19, "確定引用清單(5 筆)", lside="left")
lx = Inches(9.45); ly = Inches(1.65)
box(s, lx, ly, Inches(3.45), Inches(4.85), fill=BG, line=LGRAY)
text(s, lx + Inches(0.2), ly + Inches(0.2), Inches(3.05), Inches(0.5), [[("這一頁在做什麼", 15, True, NAVY)]])
bullets(s, [(0, "問:510.7 引用哪些外部規範", DARK, False),
            (0, "系統自動判定為關係題,走 Agent", DARK, False),
            (0, "呼叫 section_lookup 沿圖查詢", DARK, False),
            (0, "回傳結構化清單,非自由發揮", DARK, False),
            (0, "5 份規範,可驗證、不幻覺", GREEN, True)],
        l=lx + Inches(0.3), t=ly + Inches(0.75), w=Inches(2.95), sz=13, gap=11)

# 23 · 系統功能一覽 + 設定
s = newslide("PART 3 · 系統一覽", "一套系統,涵蓋匯入 → 問答 → 圖譜 → 設定全流程")
funcs = [("① 文件列表", "檢視 / 下載規範", BLUE), ("② 建立文件", "上傳 + OCR + 向量化", GREEN),
         ("③ RAG 問答", "三模式查詢", ORANGE), ("④ 知識圖譜", "引用網 / 版本鏈", PURPLE),
         ("⑤ 我的筆記本", "收藏查詢結果", TEAL), ("⑥ 管理介面", "模型 / 提示詞設定", NAVY),
         ("⑦ 向量查詢測試", "檢驗檢索品質", RGBColor(0xC2, 0x41, 0x0C)), ("⑧ 向量庫健康", "資料一致性檢查", RGBColor(0x2B, 0x8A, 0x3E))]
x0 = Inches(0.7); y0 = Inches(1.55); cw = Inches(2.92); ch = Inches(1.35); gx = Inches(0.16); gy = Inches(0.18)
for i, (h, d, col) in enumerate(funcs):
    r, c = divmod(i, 4)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    box(s, cx, cy, cw, ch, fill=WHITE, line=LGRAY)
    box(s, cx, cy, Pt(6), ch, fill=col)
    text(s, cx + Inches(0.2), cy + Inches(0.16), cw - Inches(0.35), Inches(0.4), [[(h, 14.5, True, col)]])
    text(s, cx + Inches(0.2), cy + Inches(0.68), cw - Inches(0.35), Inches(0.55), [[(d, 12.5, False, DARK)]])
place_img(s, os.path.join(DECK, "09b_settings.png"), Inches(0.7), Inches(4.75), Inches(6.0), Inches(2.3))
place_img(s, os.path.join(DECK, "09c_prompts.png"), Inches(6.9), Inches(4.75), Inches(5.7), Inches(2.3))
text(s, Inches(0.7), Inches(4.5), Inches(6.0), Inches(0.3), [[("⑥ 系統設定:模型 / OCR / 向量化參數", 12.5, True, NAVY)]], align=PP_ALIGN.CENTER)
text(s, Inches(6.9), Inches(4.5), Inches(5.7), Inches(0.3), [[("⑥ 提示詞管理:可調整 RAG 生成規則", 12.5, True, NAVY)]], align=PP_ALIGN.CENTER)

# =====================================================================
# Part 4 divider
# =====================================================================
section("PART 4", "成果與未來", "AS-IS / TO-BE · 為何適合測試規範 · 下一步")

# 24 · AS-IS / TO-BE
s = newslide("PART 4 · 成果", "AS-IS vs TO-BE:同樣的問題,兩種結果")
rows = [
    ("單篇內容題（濕度測試怎麼做）", "常漏關鍵條件、無出處", "完整條件 + 可點 PDF 頁"),
    ("換句話問（濕度 vs humidity）", "字面對不上就漏檢", "語意命中,不管怎麼問"),
    ("跨規範引用（510.7 引用哪些）", "答不全、時多時少、會幻覺", "確定清單 5 筆,可驗證"),
    ("版本追查（哪一版取代舊版）", "模型不知道", "沿版本鏈給確定答案"),
    ("未收錄規範（ASTM B117）", "硬編一段錯的", "明確說「未收錄」,不亂編"),
]
y = Inches(1.55); rh = Inches(0.86); x0 = Inches(0.55)
w1, w2, w3 = Inches(4.4), Inches(3.9), Inches(3.9)
hdr = [("問題類型", w1, NAVY), ("AS-IS｜基本 RAG", w2, RED), ("TO-BE｜RAG × KG", w3, GREEN)]
cx = x0
for h, w, col in hdr:
    box(s, cx, y, w, Inches(0.55), fill=col)
    text(s, cx, y, w, Inches(0.55), [[(h, 14.5, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i, (a, b, c) in enumerate(rows):
    ry = y + Inches(0.55) + i * rh
    cx = x0
    for j, (val, w) in enumerate(((a, w1), (b, w2), (c, w3))):
        fill = WHITE if i % 2 == 0 else BG
        if j == 2:
            fill = LGREEN
        box(s, cx, ry, w, rh, fill=fill, line=LGRAY)
        col = DARK
        bold = False
        if j == 1:
            col = RED
        if j == 2:
            col = RGBColor(0x1B, 0x6E, 0x33); bold = True
        text(s, cx + Inches(0.15), ry, w - Inches(0.3), rh, [[(val, 13, bold, col)]], anchor=MSO_ANCHOR.MIDDLE, align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
        cx += w

# 25 · 為何測試規範是最佳場景
s = newslide("PART 4 · 為什麼是測試規範", "為什麼「測試規範」是 RAG × KG 的最佳落地場景")
feats = [("📐 高度結構化", "方法 / 章節 / 程序層層分明,天生適合切塊與建節點。", BLUE),
         ("🔗 大量交叉引用", "一個方法引用 5～10 份外部規範,正是 KG 最能發揮的地方。", GREEN),
         ("🕐 版本嚴謹", "A/B/C 版取代關係明確,可沿版本鏈精準追查。", PURPLE),
         ("⚠️ 答錯代價高", "測試條件錯 → 產品驗證失效,對「可溯源」需求極強。", RED)]
x0 = Inches(0.7); y0 = Inches(1.6); cw = Inches(5.86); ch = Inches(1.95); gx = Inches(0.21); gy = Inches(0.25)
for i, (h, b, col) in enumerate(feats):
    r, c = divmod(i, 2)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    box(s, cx, cy, cw, ch, fill=WHITE, line=LGRAY, line_w=1.2)
    box(s, cx, cy, Pt(7), ch, fill=col)
    text(s, cx + Inches(0.3), cy + Inches(0.25), cw - Inches(0.5), Inches(0.5), [[(h, 18, True, col)]])
    text(s, cx + Inches(0.35), cy + Inches(0.9), cw - Inches(0.6), Inches(0.95), [[(b, 14.5, False, DARK)]])
box(s, Inches(0.7), Inches(6.35), Inches(11.93), Inches(0.7), fill=NAVY)
text(s, Inches(0.95), Inches(6.39), Inches(11.5), Inches(0.62),
     [[("金句:越是「關係密集 + 正確性要求高」的領域,知識圖譜的價值就越大 — 測試規範正中甜蜜點。", 15.5, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)

# 26 · 未來:輸入全部規範
s = newslide("PART 4 · 未來展望", "下一步之一:把「全部規範」都收進來")
text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.6),
     [[("目前已收錄 40 份軍規文件。架構已就緒,可持續匯入 MIL / ASTM / IEC / ISO / IEEE… 各類標準。", 16, False, DARK)]])
box(s, Inches(0.7), Inches(2.35), Inches(5.86), Inches(3.5), fill=LBLUE, line=BLUE, line_w=1.3)
text(s, Inches(0.95), Inches(2.5), Inches(5.4), Inches(0.5), [[("越多文件 → 越有價值（網路效應）", 16.5, True, BLUE)]])
bullets(s, [(0, "每新增一份規範,它的引用會連到既有節點", DARK, False),
            (0, "圖譜越大,能回答的跨規範關係題越多", DARK, False),
            (0, "版本鏈越完整,追查越準", DARK, False),
            (0, "價值不是線性成長,而是隨連線加乘", BLUE, True)],
        l=Inches(1.0), t=Inches(3.1), w=Inches(5.3), sz=14.5, gap=13)
box(s, Inches(6.77), Inches(2.35), Inches(5.86), Inches(3.5), fill=BG, line=LGRAY)
text(s, Inches(7.02), Inches(2.5), Inches(5.4), Inches(0.5), [[("已解決:大量匯入不再塞車", 16.5, True, GREEN)]])
bullets(s, [(0, "KG 抽取改為單一序列佇列,一份接一份", DARK, False),
            (0, "避免同時寫入造成資料庫鎖死", DARK, False),
            (0, "使用者自行陸續匯入,穩定不出錯", DARK, False),
            (0, "已壓力測試驗證,可放心擴充", GREEN, True)],
        l=Inches(7.07), t=Inches(3.1), w=Inches(5.3), sz=14.5, gap=13)
text(s, Inches(0.7), Inches(6.15), Inches(11.93), Inches(0.6),
     [[("願景:成長為一座「軍規/測試標準的活知識庫」,新進工程師也能一問就上手。", 15.5, True, NAVY)]], align=PP_ALIGN.CENTER)

# 27 · 非規範文件相容性
s = newslide("PART 4 · 未來展望", "下一步之二:對「非規範文件」也相容")
text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.6),
     [[("常見疑問:那測試報告、SOP、技術手冊、會議記錄這種「沒有規範編號」的文件,放得進來嗎?", 16, True, NAVY)]])
box(s, Inches(0.7), Inches(2.45), Inches(5.86), Inches(3.5), fill=LGREEN, line=GREEN, line_w=1.3)
text(s, Inches(0.95), Inches(2.6), Inches(5.4), Inches(0.5), [[("RAG 部分:天生相容", 17, True, GREEN)]])
bullets(s, [(0, "任何 PDF / 文字文件都能切塊、向量化", DARK, False),
            (0, "測試報告、SOP、手冊照樣能語意檢索", DARK, False),
            (0, "照樣附出處、照樣防幻覺", DARK, False),
            (0, "「內容問答」對文件類型沒有限制", GREEN, True)],
        l=Inches(1.0), t=Inches(3.2), w=Inches(5.3), sz=14.5, gap=13)
box(s, Inches(6.77), Inches(2.45), Inches(5.86), Inches(3.5), fill=LBLUE, line=BLUE, line_w=1.3)
text(s, Inches(7.02), Inches(2.6), Inches(5.4), Inches(0.5), [[("KG 部分:自動退化,不會壞", 17, True, BLUE)]])
bullets(s, [(0, "沒有規範編號 → 就不建引用節點", DARK, False),
            (0, "系統自動退回「純內容檢索」模式", DARK, False),
            (0, "有關係就多一層價值,沒關係也照常可用", DARK, False),
            (0, "規範與非規範文件可並存於同一庫", BLUE, True)],
        l=Inches(7.07), t=Inches(3.2), w=Inches(5.3), sz=14.5, gap=13)
box(s, Inches(0.7), Inches(6.15), Inches(11.93), Inches(0.75), fill=YELLOW, line=GOLD, line_w=1)
text(s, Inches(0.95), Inches(6.2), Inches(11.5), Inches(0.65),
     [[("💡 設計上「規範文件加值、非規範文件不減分」:KG 是加分項,缺了它 RAG 仍完整運作。", 15, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)

# 28 · 總結
s = newslide("總結", "一頁帶走:我們蓋了一座什麼樣的知識庫")
journey = ["基本 RAG\n查得到", "五道強化\n查得準", "＋知識圖譜\n查得到關係", "混合路由\n自動選對路", "測試規範\n活知識庫"]
n = len(journey); x = Inches(0.55); y = Inches(1.7); bw = Inches(2.25); bh = Inches(1.35); g = Inches(0.18)
cols = [GRAY, BLUE, GREEN, PURPLE, NAVY]
for i, st in enumerate(journey):
    bx = x + i * (bw + g)
    c = box(s, bx, y, bw, bh, fill=cols[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for j, ln in enumerate(st.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, 15 if j == 0 else 11.5, True, WHITE)
    if i < n - 1:
        arrow(s, bx + bw + Emu(4000), y + bh // 2 - Inches(0.1), g - Emu(8000), Inches(0.2), color=LGRAY)
pts = [
    ("查得準", "五道強化:聰明切塊、語意檢索、重排序、強制出處、防幻覺 — 解決單篇內容的準確性。", BLUE),
    ("查得到關係", "知識圖譜補上跨文件的引用與版本鏈,回答純 RAG 答不出的關係題。", GREEN),
    ("可信任", "答案可溯源、找不到會明說 — 對合規與品保場景,這是能不能用的底線。", ORANGE),
    ("可成長", "序列佇列讓匯入不塞車;規範加值、非規範相容 — 可持續長大的活知識庫。", PURPLE),
]
y0 = Inches(3.4)
for i, (h, b, col) in enumerate(pts):
    cy = y0 + i * Inches(0.85)
    box(s, Inches(0.7), cy, Inches(11.93), Inches(0.75), fill=WHITE, line=LGRAY)
    box(s, Inches(0.7), cy, Inches(2.5), Inches(0.75), fill=col)
    text(s, Inches(0.7), cy, Inches(2.5), Inches(0.75), [[(h, 16, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.4), cy, Inches(9.0), Inches(0.75), [[(b, 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)

# 29 · 結語
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(4.5), SW, Pt(4), fill=BLUE)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.2),
     [[("讓查規範,像問同事一樣簡單。", 34, True, WHITE)],
      [("而且答案有出處、有關聯、可信任。", 24, True, RGBColor(0xCF, 0xE2, 0xFF))]])
text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.6), [[("RAG × 知識圖譜　·　測試規範智慧問答系統", 16, False, GRAY)]])
text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5), [[("敬請指教 · Q & A", 18, True, RGBColor(0x9E, 0xC5, 0xFF))]])
PAGE[0] += 1

prs.save(OUT)
print("SAVED slides:", len(prs.slides._sldIdLst))
